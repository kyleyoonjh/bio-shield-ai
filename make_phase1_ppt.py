"""
OpenBioShield Phase 1 — 연구원 가이드 PPT
DNA 탐색 & 프라이머 구조 분석
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

DARK_BG  = RGBColor(0x0F, 0x17, 0x2A)
ACCENT   = RGBColor(0x10, 0xB9, 0x81)
ACCENT2  = RGBColor(0x6E, 0xE7, 0xB7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0x1E, 0x2D, 0x40)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)
BLUE     = RGBColor(0x60, 0xA5, 0xFA)
PINK     = RGBColor(0xF4, 0x72, 0x7E)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)
PURPLE   = RGBColor(0xC0, 0x84, 0xFC)
CARD2    = RGBColor(0x16, 0x22, 0x34)
TEAL     = RGBColor(0x2D, 0xD4, 0xBF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def slide():   return prs.slides.add_slide(blank)

def bg(s, color=DARK_BG):
    sh = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def box(s, x, y, w, h, fill=LIGHT_BG, radius=False):
    sh = s.shapes.add_shape(17 if radius else 1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    if radius: sh.adjustments[0] = 0.05
    return sh

def top_bar(s, color=ACCENT):
    sh = s.shapes.add_shape(1, 0, Inches(0.17), Inches(1.6), Inches(0.065))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def side_bar(s, color=ACCENT):
    sh = s.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def txt(s, text, x, y, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic

def col_bar(s, x, y, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.065), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def tip_box(s, x, y, w, content, color=YELLOW):
    box(s, x, y, w, 0.72, fill=RGBColor(0x2A, 0x20, 0x05), radius=True)
    col_bar(s, x, y, 0.72, color)
    txt(s, f"💡  {content}", x+0.18, y+0.1, w-0.25, 0.52, size=12, color=YELLOW, wrap=True)

def badge(s, x, y, label, color):
    b = s.shapes.add_shape(17, Inches(x), Inches(y), Inches(1.6), Inches(0.38))
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    b.adjustments[0] = 0.5
    txt(s, label, x+0.05, y+0.04, 1.5, 0.32, size=11, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# Slide 1 — 타이틀
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s); side_bar(s, TEAL)
txt(s, "Phase 1", 0.35, 1.1, 10, 0.55, size=16, color=TEAL, bold=True)
txt(s, "DNA 탐색 &\n변이 분석", 0.35, 1.6, 11, 1.8, size=52, bold=True)
txt(s, "COVID-19 역학 · 유전체 브라우저 · 프라이머 구조 시각화", 0.35, 3.55, 12, 0.6, size=20, color=ACCENT2)
txt(s, "시약 연구원을 위한 기능 사용 가이드", 0.35, 4.25, 10, 0.5, size=15, color=MUTED)

# 우측 기능 뱃지 3개
features = [("🦠", "COVID 대시보드"), ("🧬", "DNA 서열 맵"), ("🔬", "프라이머 구조")]
for i, (ic, lb) in enumerate(features):
    box(s, 9.6, 1.5 + i*1.55, 3.4, 1.3, fill=LIGHT_BG, radius=True)
    col_bar(s, 9.6, 1.5 + i*1.55, 1.3, TEAL)
    txt(s, ic, 9.8,  1.6  + i*1.55, 0.8, 0.9, size=30)
    txt(s, lb, 10.6, 1.72 + i*1.55, 2.2, 0.6, size=15, color=ACCENT2, bold=True)


# ════════════════════════════════════════════════════════════════
# Slide 2 — Phase 1 한눈에 보기
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s); top_bar(s, TEAL)
txt(s, "Phase 1 개요", 0.4, 0.1, 6, 0.38, size=13, color=TEAL, bold=True)
txt(s, "3개 화면으로 구성된 탐색 & 분석 워크플로우", 0.4, 0.45, 12, 0.55, size=24, bold=True)

# 3단계 흐름
flows = [
    ("1", "🦠 COVID 대시보드",    PINK,   "역학 데이터 파악",
     ["국가별·지역별 발생 현황", "치명률(CFR) 트렌드", "대한민국 시도별 현황", "타겟 바이러스 선택"]),
    ("2", "🧬 DNA 서열 맵",        BLUE,   "유전자 수준 분석",
     ["참조 서열 로드 (NCBI)", "변이 위치 오버레이", "관심 영역 줌인", "프라이머 결합 위치 확인"]),
    ("3", "🔬 프라이머 구조 분석", TEAL,  "설계 적합성 평가",
     ["2차 구조 아크 다이어그램", "GC% · Tm 수치 확인", "변이주별 미스매치 검출", "왓슨-크릭 염기쌍 시각화"]),
]
for i, (num, title, color, sub, bullets) in enumerate(flows):
    bx = box(s, 0.3 + i*4.35, 1.3, 4.1, 5.65, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.3 + i*4.35), Inches(1.3), Inches(4.1), Inches(0.08))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    # 번호 원
    circ = s.shapes.add_shape(9, Inches(0.4 + i*4.35), Inches(1.42), Inches(0.42), Inches(0.42))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    txt(s, num,   0.41 + i*4.35, 1.44, 0.4, 0.38, size=14, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, 0.9  + i*4.35, 1.43, 3.2, 0.45, size=15, bold=True, color=color)
    txt(s, sub,   0.9  + i*4.35, 1.9,  3.2, 0.38, size=12, color=MUTED)
    for j, b in enumerate(bullets):
        txt(s, f"  ›  {b}", 0.5 + i*4.35, 2.45 + j*0.72, 3.7, 0.65, size=13, color=WHITE if j == 0 else MUTED, wrap=True)
    if i < 2:
        txt(s, "➜", 4.5 + i*4.35, 4.0, 0.3, 0.5, size=22, color=color, align=PP_ALIGN.CENTER)

tip_box(s, 0.3, 7.05, 12.7,
        "Phase 1은 실험 설계 전 배경 파악용입니다. COVID 대시보드 → DNA 맵 → 구조 분석 순서로 진행하면 가장 효율적입니다.")


# ════════════════════════════════════════════════════════════════
# Slide 3 — COVID 대시보드 활용
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s); top_bar(s, PINK)
txt(s, "화면 1  |  COVID 대시보드", 0.4, 0.1, 12, 0.38, size=13, color=PINK, bold=True)
txt(s, "역학 컨텍스트 파악 — 왜 이 변이주를 타겟하는가", 0.4, 0.45, 12, 0.55, size=24, bold=True)

# 왼쪽: 화면 구성 설명
box(s, 0.3, 1.3, 5.9, 5.55, fill=LIGHT_BG, radius=True)
txt(s, "📊  화면 구성 요소", 0.5, 1.42, 5.5, 0.42, size=15, bold=True, color=ACCENT2)

sections = [
    ("세계 지도",       PINK,   "국가별 확진자·사망자 색상 히트맵\n관심 지역 클릭 → 상세 수치 표시"),
    ("대한민국 지도",   ORANGE, "시도별 발생 현황 시각화\n국내 진단 수요 예측에 활용"),
    ("핵심 지표 카드",  BLUE,   "전 세계 총 확진·사망·치명률(CFR)\n날짜별 트렌드 차트"),
    ("질환 선택기",     TEAL,   "SARS-CoV-2 외 추가 질환 전환 가능\n(HPV, STI 등 — Phase 1 고정)"),
]
for i, (name, color, desc) in enumerate(sections):
    box(s, 0.45, 1.98 + i*1.12, 5.6, 1.0, fill=CARD2, radius=True)
    col_bar(s, 0.45, 1.98 + i*1.12, 1.0, color)
    txt(s, name, 0.65, 2.04 + i*1.12, 2.2, 0.42, size=13, bold=True, color=color)
    txt(s, desc, 0.65, 2.47 + i*1.12, 5.1, 0.42, size=11.5, color=MUTED, wrap=True)

# 오른쪽: 연구원 활용 시나리오
box(s, 6.45, 1.3, 6.55, 5.55, fill=LIGHT_BG, radius=True)
txt(s, "🔬  시약 연구원 활용 시나리오", 6.65, 1.42, 6.1, 0.42, size=15, bold=True, color=ACCENT2)

scenarios = [
    ("신규 진단 키트 타겟 선정",
     PINK,
     "CFR이 높고 확산 중인 변이주를 대시보드에서\n확인 → 해당 변이주 서열로 Phase 2/3 진행"),
    ("국내 제품 출시 우선순위 결정",
     ORANGE,
     "대한민국 시도별 데이터 → 유행 지역 집중\n→ 지역 특화 변이주 타겟 프라이머 개발"),
    ("임상 검증 배경 자료 수집",
     BLUE,
     "역학 데이터 스크린샷을 ELN·보고서에 첨부\n→ 개발 필요성 근거 문서화"),
]
for i, (title, color, desc) in enumerate(scenarios):
    box(s, 6.6, 2.0 + i*1.6, 6.2, 1.48, fill=CARD2, radius=True)
    col_bar(s, 6.6, 2.0 + i*1.6, 1.48, color)
    txt(s, title, 6.8, 2.08 + i*1.6, 5.8, 0.42, size=13, bold=True, color=color)
    txt(s, desc,  6.8, 2.53 + i*1.6, 5.8, 0.88, size=12, color=MUTED, wrap=True)

tip_box(s, 0.3, 7.05, 12.7,
        "대시보드는 실시간 데이터를 반영합니다. 개발 착수 전 최신 유행 상황을 확인해 타겟 변이주를 최신화하세요.")


# ════════════════════════════════════════════════════════════════
# Slide 4 — DNA 서열 맵 활용
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s); top_bar(s, BLUE)
txt(s, "화면 2  |  DNA 서열 맵 (유전체 브라우저)", 0.4, 0.1, 12, 0.38, size=13, color=BLUE, bold=True)
txt(s, "변이 위치를 유전자 지도에서 직접 확인", 0.4, 0.45, 12, 0.55, size=24, bold=True)

# 상단: 주요 기능 4개
funcs = [
    ("📡", "NCBI 참조\n서열 자동 로드",  BLUE,   "인터넷 연결 시 최신\n참조 서열 자동 수신"),
    ("📍", "변이 위치\n어노테이션",       ORANGE, "SNP·InDel 위치에\n색상 마커 표시"),
    ("🔍", "영역 줌인\n& 탐색",           PURPLE, "관심 유전자 영역\n드래그로 확대"),
    ("🎯", "프라이머\n결합 위치",          TEAL,   "기존 알려진 프라이머\n결합 구간 표시"),
]
for i, (ic, title, color, desc) in enumerate(funcs):
    bx = box(s, 0.3 + i*3.27, 1.3, 3.1, 2.0, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.3 + i*3.27), Inches(1.3), Inches(3.1), Inches(0.07))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s, ic,    0.45 + i*3.27, 1.42, 0.7, 0.65, size=28)
    txt(s, title, 1.1  + i*3.27, 1.45, 2.1, 0.65, size=14, bold=True, color=color)
    txt(s, desc,  0.45 + i*3.27, 2.13, 2.8, 0.95, size=12, color=MUTED, wrap=True)

# 하단: 연구원 활용 가이드
box(s, 0.3, 3.55, 12.7, 3.3, fill=CARD2, radius=True)
txt(s, "🔬  시약 연구원 단계별 사용법", 0.55, 3.67, 12, 0.42, size=14, bold=True, color=ACCENT2)

steps = [
    ("1", "변이주 선택", BLUE,
     "상단 드롭다운에서 분석할 변이주(예: Omicron BA.5) 선택\n→ 해당 변이의 돌연변이 목록이 지도에 자동 표시"),
    ("2", "타겟 유전자 확인", ORANGE,
     "스파이크(S)·RdRp·N 단백질 등 진단 타겟 영역 확인\n→ 변이 밀도가 낮은 구간이 프라이머 설계에 유리"),
    ("3", "보존 영역 메모", TEAL,
     "변이 마커가 없는 구간 = 보존 영역\n→ Phase 3 FASTA 업로드 시 해당 영역 포함 서열 준비"),
    ("4", "프라이머 위치 검토", PURPLE,
     "기존 공개 프라이머 결합 위치와 신규 변이의 오버랩 확인\n→ 오버랩 시 기존 프라이머 성능 저하 가능 → 재설계 필요"),
]
for i, (num, title, color, desc) in enumerate(steps):
    circ = s.shapes.add_shape(9, Inches(0.55 + i*3.25), Inches(4.22), Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    txt(s, num,   0.56 + i*3.25, 4.24, 0.36, 0.34, size=13, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, 1.02 + i*3.25, 4.24, 2.0,  0.38, size=13, bold=True, color=color)
    txt(s, desc,  0.55 + i*3.25, 4.73, 3.0,  1.85, size=11.5, color=MUTED, wrap=True)

tip_box(s, 0.3, 7.05, 12.7,
        "변이 마커가 밀집한 구간(핫스팟)은 피하고, 여러 변이주에서 공통으로 비어있는 구간을 프라이머 타겟으로 선택하세요.")


# ════════════════════════════════════════════════════════════════
# Slide 5 — 프라이머 구조 분석 활용
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s); top_bar(s, TEAL)
txt(s, "화면 3  |  프라이머 구조 분석", 0.4, 0.1, 12, 0.38, size=13, color=TEAL, bold=True)
txt(s, "기존 프라이머의 안정성 · 변이 영향 사전 평가", 0.4, 0.45, 12, 0.55, size=24, bold=True)

# 왼쪽: 시각화 설명
box(s, 0.3, 1.3, 5.9, 5.55, fill=LIGHT_BG, radius=True)
txt(s, "📐  화면 구성 요소", 0.5, 1.42, 5.5, 0.42, size=15, bold=True, color=ACCENT2)

viz_items = [
    ("아크 다이어그램", TEAL,
     "프라이머 2차 구조를 반원 아크로 표현\n아크 = 왓슨-크릭 염기쌍 결합\n꼬임·겹침이 많으면 헤어핀 위험"),
    ("염기 조성 막대", BLUE,
     "A·T·G·C 각 염기의 위치별 비율 표시\n GC 편중 구간을 시각적으로 파악"),
    ("열역학 지표 카드", ORANGE,
     "GC% · Tm · MFE (최소 자유 에너지)\n수치가 낮을수록 구조 형성 위험 높음"),
    ("변이주 미스매치", PINK,
     "선택 변이주에서 해당 프라이머의\n미스매치 위치·개수 빨간 마커 표시"),
]
for i, (title, color, desc) in enumerate(viz_items):
    box(s, 0.45, 2.02 + i*1.12, 5.6, 1.0, fill=CARD2, radius=True)
    col_bar(s, 0.45, 2.02 + i*1.12, 1.0, color)
    txt(s, title, 0.65, 2.08 + i*1.12, 2.5, 0.42, size=13, bold=True, color=color)
    txt(s, desc,  0.65, 2.52 + i*1.12, 5.1, 0.42, size=11.5, color=MUTED, wrap=True)

# 오른쪽: 수치 해석 가이드
box(s, 6.45, 1.3, 6.55, 2.55, fill=LIGHT_BG, radius=True)
txt(s, "📊  수치 판독 기준", 6.65, 1.42, 6.1, 0.42, size=15, bold=True, color=ACCENT2)

metrics_guide = [
    ("GC%",   "40–60%",         "이상 범위",      TEAL),
    ("Tm",    "58–65°C",        "qPCR 권장",      BLUE),
    ("MFE",   "> -6 kcal/mol",  "구조 위험 없음", ORANGE),
    ("미스매치", "0–1개",        "2개↑ 재설계",   PINK),
]
for i, (metric, ideal, label, color) in enumerate(metrics_guide):
    box(s, 6.6, 1.95 + i*0.52, 6.25, 0.46, fill=CARD2, radius=True)
    txt(s, metric, 6.75, 2.0  + i*0.52, 1.8, 0.38, size=12, bold=True, color=color)
    txt(s, ideal,  8.6,  2.0  + i*0.52, 2.0, 0.38, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label,  10.6, 2.0  + i*0.52, 2.1, 0.38, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# 오른쪽 하단: 활용 시나리오
box(s, 6.45, 4.05, 6.55, 2.8, fill=LIGHT_BG, radius=True)
txt(s, "🔬  실무 활용 케이스", 6.65, 4.17, 6.1, 0.42, size=15, bold=True, color=ACCENT2)

cases = [
    ("기존 키트 변이 영향 평가", TEAL,
     "현재 시판 프라이머 서열 입력 → 신규 변이주 선택\n→ 미스매치 개수 확인 → 2개↑이면 재설계 착수"),
    ("Phase 3 결과 교차 검증", ORANGE,
     "Phase 3에서 설계된 신규 프라이머를\n이 화면에서 2차 구조·Tm 재확인"),
]
for i, (title, color, desc) in enumerate(cases):
    box(s, 6.6, 4.65 + i*1.0, 6.25, 0.88, fill=CARD2, radius=True)
    col_bar(s, 6.6, 4.65 + i*1.0, 0.88, color)
    txt(s, title, 6.8, 4.72 + i*1.0, 5.8, 0.38, size=12, bold=True, color=color)
    txt(s, desc,  6.8, 5.12 + i*1.0, 5.8, 0.38, size=11.5, color=MUTED, wrap=True)

tip_box(s, 0.3, 7.05, 12.7,
        "MFE 값이 -6 kcal/mol보다 낮으면(더 안정적이면) 헤어핀 구조 가능성이 높습니다. 해당 프라이머는 Phase 3 재설계를 권장합니다.")


# ════════════════════════════════════════════════════════════════
# Slide 6 — 변이주 선택 & 질환 전환
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s); top_bar(s, PURPLE)
txt(s, "공통 기능  |  변이주 선택 & 질환 전환", 0.4, 0.1, 12, 0.38, size=13, color=PURPLE, bold=True)
txt(s, "분석 대상을 정확하게 설정하는 방법", 0.4, 0.45, 12, 0.55, size=24, bold=True)

# 변이주 목록 (왼쪽)
box(s, 0.3, 1.3, 6.1, 5.55, fill=LIGHT_BG, radius=True)
txt(s, "🦠  지원 변이주 목록 (SARS-CoV-2 기준)", 0.5, 1.42, 5.8, 0.42, size=14, bold=True, color=ACCENT2)

variants = [
    ("Wild-type",   "초기 우한 원형 서열",          MUTED),
    ("Alpha",       "B.1.1.7 — 영국발, D614G+N501Y",   BLUE),
    ("Beta",        "B.1.351 — 남아공발, E484K",         TEAL),
    ("Delta",       "B.1.617.2 — 인도발, L452R",         YELLOW),
    ("Omicron BA.1","돌연변이 32개 — 기존 프라이머 영향 큼", PINK),
    ("Omicron BA.5","국내 주요 유행 아형",            ORANGE),
    ("XBB / JN.1",  "최신 재조합 변이주",              PURPLE),
]
for i, (name, desc, color) in enumerate(variants):
    box(s, 0.45, 1.98 + i*0.63, 5.8, 0.56, fill=CARD2, radius=True)
    col_bar(s, 0.45, 1.98 + i*0.63, 0.56, color)
    txt(s, name, 0.65, 2.03 + i*0.63, 2.0, 0.46, size=12, bold=True, color=color)
    txt(s, desc, 2.7,  2.03 + i*0.63, 3.4, 0.46, size=11.5, color=MUTED)

# 질환 전환 (오른쪽 상단)
box(s, 6.6, 1.3, 6.4, 2.5, fill=LIGHT_BG, radius=True)
txt(s, "🔄  지원 질환 전환", 6.8, 1.42, 5.9, 0.42, size=14, bold=True, color=ACCENT2)
diseases = [("SARS-CoV-2", PINK, "COVID-19 진단"),
            ("HPV",        ORANGE, "자궁경부암 관련"),
            ("STI",        PURPLE, "성매개 감염 패널")]
for i, (d, c, sub) in enumerate(diseases):
    box(s, 6.75, 1.98 + i*0.62, 6.1, 0.54, fill=CARD2, radius=True)
    col_bar(s, 6.75, 1.98 + i*0.62, 0.54, c)
    txt(s, d,   6.95, 2.03 + i*0.62, 2.0, 0.44, size=13, bold=True, color=c)
    txt(s, sub, 9.0,  2.03 + i*0.62, 3.6, 0.44, size=12, color=MUTED)

# 연구원 가이드 (오른쪽 하단)
box(s, 6.6, 4.0, 6.4, 2.85, fill=LIGHT_BG, radius=True)
txt(s, "📌  변이주 선택 실무 팁", 6.8, 4.12, 5.9, 0.42, size=14, bold=True, color=ACCENT2)
tips = [
    "현재 국내 유행 아형부터 분석 시작",
    "기존 키트 평가 시 → 출시 당시 유행 변이 + 최신 변이 두 가지 비교",
    "변이주별 미스매치가 다르면 → 가장 많은 변이주에 Omicron 계열 선택",
    "신규 변이 추가는 개발팀 요청으로 반영 가능",
]
for i, t in enumerate(tips):
    txt(s, f"  ›  {t}", 6.8, 4.65 + i*0.52, 5.9, 0.48, size=12, color=WHITE if i == 0 else MUTED, wrap=True)

tip_box(s, 0.3, 7.05, 12.7,
        "Omicron BA.1은 스파이크 영역에 돌연변이가 32개로 가장 많습니다. 기존 스파이크 타겟 프라이머를 보유 중이라면 반드시 재검토하세요.")


# ════════════════════════════════════════════════════════════════
# Slide 7 — Phase 1 → Phase 2/3 연결
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s); top_bar(s, ACCENT)
txt(s, "Phase 1 활용 결론", 0.4, 0.1, 12, 0.38, size=13, color=ACCENT, bold=True)
txt(s, "Phase 1에서 얻은 정보를 다음 단계로 연결하기", 0.4, 0.45, 12, 0.55, size=24, bold=True)

box(s, 0.3, 1.25, 12.7, 5.55, fill=LIGHT_BG, radius=True)

connections = [
    ("Phase 1에서 얻는 것", TEAL, [
        "타겟 변이주 결정 (대시보드 역학 데이터 기반)",
        "프라이머 결합 후보 영역 (DNA 맵 보존 구간)",
        "기존 프라이머 재설계 필요 여부 (구조 분석 미스매치)",
        "개발 배경 근거 자료 (역학 현황 스크린샷)",
    ]),
    ("→  Phase 2에 전달", YELLOW, [
        "기존 프라이머 서열 → EP05/EP09 정밀도 검증",
        "변이 위치 정보 → PFPS 위험도 계산 입력값",
        "미스매치 위치 → 3'-end 근접 여부 확인 (HIGH 위험)",
        "대시보드 역학 데이터 → 검증 우선순위 결정",
    ]),
    ("→  Phase 3에 전달", ACCENT, [
        "보존 영역 좌표 → FASTA 파일 구성 기준",
        "타겟 변이주 서열 → FASTA에 포함할 변이주 목록",
        "어세이 타입 결정 → qPCR / Multiplex 선택 근거",
        "Tm 기준 파악 → Advanced Options 설정값 참고",
    ]),
]
for i, (title, color, items) in enumerate(connections):
    col_x = 0.5 + i*4.2
    txt(s, title, col_x, 1.45, 3.9, 0.5, size=14, bold=True, color=color)
    for j, item in enumerate(items):
        box(s, col_x, 2.05 + j*0.9, 4.0, 0.78, fill=CARD2, radius=True)
        col_bar(s, col_x, 2.05 + j*0.9, 0.78, color)
        txt(s, item, col_x + 0.2, 2.12 + j*0.9, 3.75, 0.62, size=12, color=WHITE if j == 0 else MUTED, wrap=True)

tip_box(s, 0.3, 7.05, 12.7,
        "Phase 1은 분석 방향을 결정하는 단계입니다. 탐색 없이 바로 Phase 3으로 넘어가면 FASTA 구성 실수로 재작업이 생길 수 있습니다.")


out = os.path.join(os.path.dirname(__file__), "OpenBioShield_Phase1_연구원가이드.pptx")
prs.save(out)
print(f"저장 완료: {out}")
