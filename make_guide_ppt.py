"""
OpenBioShield 기능 설명 PPT — 시약 연구원 가이드
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 색상 팔레트 ──────────────────────────────────────────────
DARK_BG  = RGBColor(0x0F, 0x17, 0x2A)
ACCENT   = RGBColor(0x10, 0xB9, 0x81)
ACCENT2  = RGBColor(0x6E, 0xE7, 0xB7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0x1E, 0x2D, 0x40)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)
BLUE     = RGBColor(0x60, 0xA5, 0xFA)
PINK     = RGBColor(0xF4, 0x72, 0x7E)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)
PURPLE   = RGBColor(0xC0, 0x84, 0xFC)
CARD2    = RGBColor(0x16, 0x22, 0x34)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(blank)

def bg(s, color=DARK_BG):
    sh = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def box(s, x, y, w, h, fill=LIGHT_BG, radius=False):
    sh = s.shapes.add_shape(17 if radius else 1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    if radius: sh.adjustments[0] = 0.05
    return sh

def top_bar(s, color=ACCENT, h=0.07):
    sh = s.shapes.add_shape(1, 0, Inches(0.17), Inches(1.4), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def txt(s, text, x, y, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic

def step_badge(s, x, y, num, color=ACCENT):
    circ = s.shapes.add_shape(9, Inches(x), Inches(y), Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    txt(s, str(num), x+0.01, y+0.01, 0.36, 0.36, size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

def tip_box(s, x, y, w, content, color=YELLOW):
    box(s, x, y, w, 0.75, fill=RGBColor(0x2A, 0x20, 0x05), radius=True)
    lb = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(0.75))
    lb.fill.solid(); lb.fill.fore_color.rgb = color; lb.line.fill.background()
    txt(s, f"💡  {content}", x+0.15, y+0.1, w-0.2, 0.55, size=12, color=YELLOW)

def scenario_box(s, x, y, w, h, title, content, color=BLUE):
    box(s, x, y, w, h, fill=RGBColor(0x05, 0x10, 0x2A), radius=True)
    lb = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(h))
    lb.fill.solid(); lb.fill.fore_color.rgb = color; lb.line.fill.background()
    txt(s, title,   x+0.2, y+0.12, w-0.3, 0.4,    size=12, bold=True, color=color)
    txt(s, content, x+0.2, y+0.55, w-0.3, h-0.65, size=11.5, color=MUTED, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 1 — 타이틀
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
bar = s.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

txt(s, "OpenBioShield", 0.35, 1.3, 12, 1.1, size=50, bold=True)
txt(s, "시약 연구원을 위한 기능 사용 가이드", 0.35, 2.55, 12, 0.7, size=26, color=ACCENT2)
txt(s, "각 단계별 실무 활용법 · 주요 수치 해석 · 실험 전 체크포인트", 0.35, 3.35, 12, 0.55, size=16, color=MUTED)

# 우측 아이콘 박스들
icons = [("🧬", "서열 입력"), ("⚙️", "자동 설계"), ("📊", "결과 분석"), ("📄", "보고서")]
for i, (ic, lb) in enumerate(icons):
    box(s, 9.5, 1.5 + i*1.3, 3.2, 1.1, fill=LIGHT_BG, radius=True)
    txt(s, ic, 9.65, 1.55 + i*1.3, 0.7, 0.9, size=26)
    txt(s, lb, 10.4, 1.7  + i*1.3, 2.2, 0.6, size=14, color=ACCENT2)

txt(s, "총 9단계 자동화 파이프라인", 0.35, 4.2, 8, 0.45, size=14, color=MUTED)


# ════════════════════════════════════════════════════════════════
# Slide 2 — 워크플로우 한눈에 보기
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s)
txt(s, "전체 워크플로우",  0.4, 0.1, 5, 0.4, size=13, color=ACCENT, bold=True)
txt(s, "FASTA 업로드부터 보고서까지 — 한눈에 보기", 0.4, 0.48, 12, 0.6, size=24, bold=True)

steps = [
    ("1", "FASTA\n업로드",    ACCENT,  "서열 파일\n입력"),
    ("2", "MSA\n정렬",        BLUE,    "MAFFT\n자동 정렬"),
    ("3", "보존 영역\n탐색",  BLUE,    "Shannon\nEntropy"),
    ("4", "후보\n생성",       PURPLE,  "Primer3\n자동 설계"),
    ("5", "특이성\n검증",     ORANGE,  "Off-target\n필터"),
    ("6", "커버리지\n분석",   YELLOW,  "변이 대응\n계산"),
    ("7", "열역학\n분석",     PINK,    "Tm·GC·\n이합체"),
    ("8", "AI 점수\n예측",    ACCENT,  "효율\n예측"),
    ("9", "보고서\n생성",     ACCENT2, "HTML·PDF\n저장"),
]

for i, (num, label, color, sub) in enumerate(steps):
    x = 0.3 + i * 1.44
    box(s, x, 1.4, 1.28, 2.2, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(x), Inches(1.4), Inches(1.28), Inches(0.07))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    circ = s.shapes.add_shape(9, Inches(x+0.45), Inches(1.5), Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    txt(s, num,   x+0.46, 1.52, 0.36, 0.35, size=14, bold=True, align=PP_ALIGN.CENTER)
    txt(s, label, x+0.04, 2.0,  1.2,  0.65, size=11, bold=True, align=PP_ALIGN.CENTER)
    txt(s, sub,   x+0.04, 2.7,  1.2,  0.65, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    if i < 8:
        txt(s, "→", x+1.28, 2.1, 0.18, 0.4, size=14, color=MUTED, align=PP_ALIGN.CENTER)

# 연구원 액션 포인트
box(s, 0.3, 3.85, 12.7, 3.0, fill=CARD2, radius=True)
txt(s, "🔬  시약 연구원이 집중해야 할 단계", 0.55, 3.97, 12, 0.45, size=14, bold=True, color=ACCENT2)

actions = [
    (ACCENT,  "Step 1",  "FASTA 파일 품질 확인 — 서열 수가 많을수록 커버리지 정확도 향상"),
    (YELLOW,  "Step 6",  "커버리지 점수 확인 — 80% 이상인 후보만 실험 대상으로 선정 권장"),
    (PINK,    "Step 7",  "Tm 균형 확인 — 포워드/리버스 Tm 차이 ±2°C 이내 후보 우선 선택"),
    (ACCENT,  "Step 8",  "AI 점수 근거 확인 — 낮은 점수 원인을 파악해 설계 파라미터 조정"),
]
for i, (color, step, desc) in enumerate(actions):
    step_badge(s, 0.55, 4.6 + i*0.55, step, color)
    txt(s, desc, 1.05, 4.6 + i*0.55, 11.8, 0.5, size=12.5, color=WHITE)


# ════════════════════════════════════════════════════════════════
# Slide 3 — FASTA 업로드 & 파라미터 설정
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s, ACCENT)
txt(s, "Step 1  |  FASTA 업로드 & 파라미터 설정", 0.4, 0.1, 12, 0.4, size=13, color=ACCENT, bold=True)
txt(s, "좋은 입력이 좋은 결과를 만든다", 0.4, 0.48, 12, 0.6, size=24, bold=True)

# 왼쪽: FASTA 파일 요건
box(s, 0.3, 1.3, 5.8, 4.5, fill=LIGHT_BG, radius=True)
txt(s, "📁  FASTA 파일 요건", 0.5, 1.42, 5.4, 0.45, size=15, bold=True, color=ACCENT2)

fasta_tips = [
    ("✅", "권장 서열 수", "20개 이상 — 커버리지 통계 신뢰도 확보"),
    ("✅", "서열 길이",    "분석 타겟 유전자 전체 포함 (트리밍 불필요)"),
    ("✅", "헤더 형식",    ">샘플명 또는 >accession_번호"),
    ("⚠️", "주의",        "중복 서열 포함 시 커버리지 과대 계산됨"),
    ("⚠️", "주의",        "갭(-) 문자 자동 제거 처리됨"),
]
for i, (ic, key, val) in enumerate(fasta_tips):
    box(s, 0.45, 1.95 + i*0.68, 5.5, 0.62, fill=CARD2, radius=True)
    txt(s, ic,  0.55, 2.0 + i*0.68, 0.4, 0.52, size=15)
    txt(s, key, 0.95, 2.02 + i*0.68, 1.5, 0.52, size=12, bold=True, color=ACCENT2)
    txt(s, val, 2.5,  2.02 + i*0.68, 3.3, 0.52, size=11.5, color=MUTED)

# 오른쪽: 어세이 타입 선택
box(s, 6.35, 1.3, 6.65, 4.5, fill=LIGHT_BG, radius=True)
txt(s, "🔬  어세이 타입 선택 가이드", 6.55, 1.42, 6.2, 0.45, size=15, bold=True, color=ACCENT2)

assay_types = [
    ("qPCR", ACCENT,
     "정량 분석 · 형광 탐침 불필요\n제품 크기 80–200 bp 자동 최적화"),
    ("Multiplex qPCR", YELLOW,
     "TaqMan 프로브 자동 설계 포함\n여러 타겟 동시 검출 시 사용"),
    ("Standard PCR", BLUE,
     "겔 전기영동 확인용\n제품 크기 제한 완화"),
]
for i, (atype, color, desc) in enumerate(assay_types):
    box(s, 6.5, 2.0 + i*1.2, 6.3, 1.08, fill=CARD2, radius=True)
    cb = s.shapes.add_shape(1, Inches(6.5), Inches(2.0 + i*1.2), Inches(0.06), Inches(1.08))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s, atype, 6.7, 2.07 + i*1.2, 5.9, 0.4,  size=14, bold=True, color=color)
    txt(s, desc,  6.7, 2.5  + i*1.2, 5.9, 0.55, size=11.5, color=MUTED)

tip_box(s, 0.3, 6.0, 12.7,
        "동일 타겟의 변이주(variants) FASTA를 함께 포함하면 커버리지 분석의 신뢰도가 크게 올라갑니다.")


# ════════════════════════════════════════════════════════════════
# Slide 4 — 결과 화면: 후보 목록 해석
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s, PURPLE)
txt(s, "Step 4–8  |  후보 프라이머 결과 목록", 0.4, 0.1, 12, 0.4, size=13, color=PURPLE, bold=True)
txt(s, "각 컬럼이 무엇을 의미하는가", 0.4, 0.48, 12, 0.6, size=24, bold=True)

# 가상 결과 테이블 헤더
headers = ["순위", "Forward (5'→3')", "Reverse (5'→3')", "Tm Fwd", "Tm Rev", "GC%", "커버리지", "AI점수", "최종점수"]
col_w =   [0.55,   2.8,              2.8,              0.7,     0.7,     0.6,  0.85,    0.75,    0.85]
col_x = [0.3]
for w in col_w[:-1]: col_x.append(col_x[-1] + w + 0.05)

for ci, (h, cw, cx) in enumerate(zip(headers, col_w, col_x)):
    box(s, cx, 1.3, cw, 0.45, fill=ACCENT, radius=False)
    txt(s, h, cx+0.04, 1.33, cw-0.05, 0.4, size=10, bold=True, align=PP_ALIGN.CENTER)

rows_data = [
    ["1", "ATGGCTAGCTAGCTAGC…", "CGATCGATCGATCG…", "62.3°", "61.8°", "52%", "94.5%", "78.4", "88.2"],
    ["2", "GCTAGCTAGCTAGCTAG…", "TAGCTAGCTAGCTA…", "60.1°", "61.2°", "48%", "87.3%", "71.2", "82.5"],
    ["3", "CGATCGATCGATCGATC…", "ATCGATCGATCGAT…", "63.5°", "60.2°", "55%", "79.1%", "65.8", "75.3"],
]
row_fills = [RGBColor(0x05, 0x2A, 0x1A), CARD2, LIGHT_BG]
for ri, row in enumerate(rows_data):
    for ci, (cell, cw, cx) in enumerate(zip(row, col_w, col_x)):
        box(s, cx, 1.8 + ri*0.52, cw, 0.5, fill=row_fills[ri])
        color = ACCENT if (ri == 0 and ci in [6,7,8]) else (ACCENT2 if ri == 0 else MUTED)
        txt(s, cell, cx+0.04, 1.83 + ri*0.52, cw-0.05, 0.44,
            size=10, color=color, align=PP_ALIGN.CENTER, wrap=False)

# 컬럼 설명
box(s, 0.3, 3.55, 12.7, 3.55, fill=CARD2, radius=True)
txt(s, "📖  각 수치 해석 가이드", 0.55, 3.67, 12, 0.4, size=14, bold=True, color=ACCENT2)

col_guides = [
    (ACCENT, "Tm (융해온도)",
     "포워드·리버스 차이 ±2°C 이내 권장\n차이가 클수록 PCR 최적화 어려움"),
    (YELLOW, "GC%",
     "40–60% 범위가 이상적\n극단값은 비특이 증폭 위험"),
    (BLUE, "커버리지",
     "입력 서열 중 증폭 가능 비율\n80% 이상 후보를 우선 선택"),
    (PINK, "AI 점수",
     "설계 품질 종합 예측 (0–100)\n점수 낮으면 근거 클릭해 원인 확인"),
    (ACCENT2, "최종 점수",
     "커버리지·Tm·AI·특이성 가중 합산\n실험 우선순위 결정 기준"),
]
for i, (color, title, desc) in enumerate(col_guides):
    box(s, 0.4 + i*2.53, 4.2, 2.4, 2.65, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.4 + i*2.53), Inches(4.2), Inches(2.4), Inches(0.07))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s, title, 0.55 + i*2.53, 4.32, 2.1, 0.45, size=13, bold=True, color=color)
    txt(s, desc,  0.55 + i*2.53, 4.82, 2.2, 1.85, size=11.5, color=MUTED, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 5 — 커버리지 분석 활용
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s, YELLOW)
txt(s, "Step 6  |  변이 커버리지 분석", 0.4, 0.1, 12, 0.4, size=13, color=YELLOW, bold=True)
txt(s, "이 프라이머가 몇 %의 변이주를 검출할 수 있는가", 0.4, 0.48, 12, 0.6, size=24, bold=True)

# 왼쪽: 설명
box(s, 0.3, 1.3, 5.8, 5.55, fill=LIGHT_BG, radius=True)
txt(s, "커버리지란?", 0.5, 1.42, 5.4, 0.45, size=15, bold=True, color=YELLOW)
txt(s,
    "입력한 FASTA 서열 전체 중, 해당 프라이머 쌍이\n"
    "≤2개 미스매치 조건으로 결합 가능한 서열의 비율입니다.\n\n"
    "예: 서열 100개 중 94개에 결합 가능 → 커버리지 94%",
    0.5, 1.95, 5.5, 1.4, size=13, color=WHITE)

# 판정 기준 바
criteria = [
    ("90% 이상", "최우선 선택",      ACCENT,  RGBColor(0x05, 0x2A, 0x1A)),
    ("70–90%",  "조건부 선택",       YELLOW,  RGBColor(0x2A, 0x20, 0x05)),
    ("50–70%",  "파라미터 재조정",   ORANGE,  RGBColor(0x2A, 0x15, 0x05)),
    ("50% 미만", "사용 비권장",      PINK,    RGBColor(0x2A, 0x05, 0x0A)),
]
for i, (range_txt, verdict, color, fill) in enumerate(criteria):
    box(s, 0.45, 3.55 + i*0.62, 5.5, 0.56, fill=fill, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.45), Inches(3.55 + i*0.62), Inches(0.06), Inches(0.56))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s, range_txt, 0.65, 3.6 + i*0.62, 1.6, 0.46, size=13, bold=True, color=color)
    txt(s, verdict,   2.3,  3.6 + i*0.62, 3.5, 0.46, size=12.5, color=WHITE)

# 오른쪽: 실제 활용 시나리오
box(s, 6.35, 1.3, 6.65, 5.55, fill=LIGHT_BG, radius=True)
txt(s, "🔬  실무 시나리오", 6.55, 1.42, 6.2, 0.45, size=15, bold=True, color=ACCENT2)

scenarios = [
    ("SARS-CoV-2 변이주 대응",
     "Alpha·Delta·Omicron 서열을 모두 포함한\nFASTA 업로드 → 커버리지 90%+ 후보 선택\n→ 범용 진단 키트 개발 가능",
     BLUE),
    ("인플루엔자 A/B 구분",
     "A형·B형 서열 혼합 업로드 시\n커버리지가 낮은 후보 = 특정 아형 특이적\n→ 의도적으로 낮은 커버리지 활용 가능",
     PURPLE),
    ("커버리지가 낮을 때",
     "Advanced Options에서 최대 미스매치 수 조정\n(기본값 2 → 3으로 완화)\n단, 특이성 점수 동시 확인 필수",
     ORANGE),
]
for i, (title, desc, color) in enumerate(scenarios):
    scenario_box(s, 6.5, 2.0 + i*1.6, 6.3, 1.48, title, desc, color)

tip_box(s, 0.3, 7.0, 12.7,
        "커버리지 80% 미만이면 Advanced Options → Max Mismatches를 3으로 높이거나, 보존 영역 임계값을 낮춰보세요.")


# ════════════════════════════════════════════════════════════════
# Slide 6 — 열역학 점수 해석 (Tm, GC, 이합체)
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s, PINK)
txt(s, "Step 7  |  열역학 분석", 0.4, 0.1, 12, 0.4, size=13, color=PINK, bold=True)
txt(s, "Tm · GC% · 이합체 · 헤어핀 — 실험 전 검증", 0.4, 0.48, 12, 0.6, size=24, bold=True)

params = [
    ("🌡️", "융해온도 (Tm)", PINK,
     ["qPCR 권장: 58–65°C",
      "포워드·리버스 차이 ≤2°C",
      "Multiplex: 프로브 Tm이 프라이머보다 5–10°C 높아야 함"],
     "Tm 균형이 맞지 않으면 한 프라이머가 먼저 결합\n해제되어 증폭 효율이 급격히 떨어집니다."),
    ("🧮", "GC 함량", BLUE,
     ["이상 범위: 40–60%",
      "3'-end 마지막 5nt: GC 2–3개 권장",
      "GC 클램프 과강(>80%) = 비특이 결합 위험"],
     "3'-end GC 클램프는 프라이머 신장 시작 안정성을\n결정합니다. GC로 끝나는 구조가 이상적입니다."),
    ("🔗", "이합체·헤어핀", ORANGE,
     ["Δ G > -6 kcal/mol 권장 (덜 안정적일수록 좋음)",
      "3'-end 이합체는 특히 위험 → 즉각 폐기",
      "헤어핀 Tm이 어닐링 온도보다 낮아야 함"],
     "이합체 형성은 실제 PCR에서 가짜 밴드(artifact)를\n만드는 주원인입니다."),
]

for i, (icon, title, color, bullets, note) in enumerate(params):
    bx = box(s, 0.3 + i*4.35, 1.35, 4.1, 5.55, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.3 + i*4.35), Inches(1.35), Inches(4.1), Inches(0.08))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s, icon,  0.45 + i*4.35, 1.5,  0.7, 0.65, size=26)
    txt(s, title, 1.1  + i*4.35, 1.52, 3.2, 0.55, size=16, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, f"• {b}", 0.5 + i*4.35, 2.25 + j*0.7, 3.7, 0.65,
            size=12, color=WHITE if j == 0 else MUTED, wrap=True)
    # 노트 박스
    note_bg = s.shapes.add_shape(1,
        Inches(0.35 + i*4.35), Inches(4.55), Inches(4.0), Inches(1.25))
    note_bg.fill.solid(); note_bg.fill.fore_color.rgb = CARD2; note_bg.line.fill.background()
    txt(s, "💬 " + note, 0.5 + i*4.35, 4.62, 3.8, 1.1, size=11, color=MUTED, italic=True, wrap=True)

tip_box(s, 0.3, 6.05, 12.7,
        "Tm 차이가 3°C를 넘는 후보는 Gradient PCR로 최적 어닐링 온도를 찾거나 다음 순위 후보로 교체하세요.")


# ════════════════════════════════════════════════════════════════
# Slide 7 — AI 점수 & 근거 해석
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s, ACCENT)
txt(s, "Step 8  |  AI 효율 점수", 0.4, 0.1, 12, 0.4, size=13, color=ACCENT, bold=True)
txt(s, "점수가 낮다면 — 어디서 깎였는지 확인하라", 0.4, 0.48, 12, 0.6, size=24, bold=True)

# 왼쪽: 점수 구성 설명
box(s, 0.3, 1.3, 5.9, 5.55, fill=LIGHT_BG, radius=True)
txt(s, "AI 점수 구성 (0–100점)", 0.5, 1.42, 5.6, 0.45, size=15, bold=True, color=ACCENT2)

components = [
    ("기본 점수",        "+60점",  "모든 후보의 출발점"),
    ("GC% (양방향)",    "최대-16점", "40–60% 벗어나면 감점"),
    ("Tm 균형",          "최대-15점", "Tm 차이 1°C당 -3점"),
    ("3'-end GC 클램프", "최대-13점", "너무 약하면 -8, 강하면 -5"),
    ("이중가닥 엔트로피", "±보정",    "서열 복잡도 보너스"),
    ("제품 크기",        "±5점",    "80–200bp → +5, >250bp → -5"),
    ("보존 영역 품질",   "최대-10점", "엔트로피 낮을수록 유리"),
]
for i, (comp, delta, note) in enumerate(components):
    box(s, 0.45, 2.0 + i*0.6, 5.6, 0.55, fill=CARD2, radius=True)
    color = ACCENT if "+" in delta and "-" not in delta else (PINK if "-" in delta else YELLOW)
    txt(s, comp,  0.6,  2.05 + i*0.6, 2.2, 0.45, size=12, color=WHITE)
    txt(s, delta, 2.85, 2.05 + i*0.6, 1.1, 0.45, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, note,  4.0,  2.05 + i*0.6, 2.0, 0.45, size=10.5, color=MUTED)

# 오른쪽: 보고서에서 근거 확인 방법
box(s, 6.45, 1.3, 6.55, 2.7, fill=LIGHT_BG, radius=True)
txt(s, "📋  보고서에서 근거 확인", 6.65, 1.42, 6.2, 0.45, size=15, bold=True, color=ACCENT2)
txt(s,
    "HTML 보고서 하단\n\"Candidate Details — AI Score Breakdown\"\n섹션을 클릭하면 후보별 점수 세부 내역 확인 가능",
    6.65, 1.95, 6.2, 1.0, size=13, color=WHITE)

# 점수 구간별 판정
box(s, 6.45, 4.2, 6.55, 2.65, fill=LIGHT_BG, radius=True)
txt(s, "📊  점수 구간별 실험 판단", 6.65, 4.32, 6.2, 0.45, size=15, bold=True, color=ACCENT2)
judgments = [
    ("80점 이상", "바로 주문·합성 진행",       ACCENT),
    ("60–80점",  "보완 파라미터 확인 후 진행",  YELLOW),
    ("40–60점",  "설계 재검토 권장",            ORANGE),
    ("40점 미만", "후보 제외 또는 재설계",      PINK),
]
for i, (score, action, color) in enumerate(judgments):
    box(s, 6.6, 4.85 + i*0.48, 6.25, 0.43, fill=CARD2, radius=True)
    cb2 = s.shapes.add_shape(1, Inches(6.6), Inches(4.85 + i*0.48), Inches(0.05), Inches(0.43))
    cb2.fill.solid(); cb2.fill.fore_color.rgb = color; cb2.line.fill.background()
    txt(s, score,  6.75, 4.89 + i*0.48, 1.5, 0.38, size=12, bold=True, color=color)
    txt(s, action, 8.3,  4.89 + i*0.48, 4.4, 0.38, size=11.5, color=MUTED)

tip_box(s, 0.3, 7.0, 12.7,
        "AI 점수가 낮아도 커버리지·Tm이 우수하면 실험적으로 성공하는 경우가 있습니다. 최종 점수(가중합)를 기준으로 판단하세요.")


# ════════════════════════════════════════════════════════════════
# Slide 8 — 보고서 활용법
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s, BLUE)
txt(s, "Step 9  |  보고서 활용", 0.4, 0.1, 12, 0.4, size=13, color=BLUE, bold=True)
txt(s, "HTML 보고서를 최대한 활용하는 방법", 0.4, 0.48, 12, 0.6, size=24, bold=True)

sections = [
    ("📈 요약 테이블",
     ACCENT,
     ["상위 후보 전체 수치 한눈에 비교",
      "컬럼 헤더 클릭 → 정렬 기능",
      "최종 점수 기준으로 우선순위 확인"],
     "가장 먼저 확인\n→ 상위 3개 후보 선정"),
    ("🧬 서열 정보",
     YELLOW,
     ["Forward·Reverse 전체 서열 복사",
      "제품 크기, 위치 정보",
      "프로브 서열 (Multiplex 시)"],
     "합성 발주 시\n서열 그대로 복사"),
    ("🤖 AI 근거 세부",
     PURPLE,
     ["후보 클릭 → 점수 성분별 분해",
      "어느 특성이 점수 낮춘지 확인",
      "재설계 방향 파악 가능"],
     "점수 낮은 이유\n파악할 때"),
    ("📥 PDF 다운로드",
     PINK,
     ["보고서 전체 PDF 저장",
      "팀 공유·아카이빙 용도",
      "ELN(전자실험노트) 첨부"],
     "연구 기록 보관\n· 팀 공유"),
]

for i, (title, color, items, use_case) in enumerate(sections):
    bx = box(s, 0.3 + i*3.27, 1.35, 3.1, 5.55, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.3 + i*3.27), Inches(1.35), Inches(3.1), Inches(0.08))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s, title, 0.45 + i*3.27, 1.5,  2.85, 0.5, size=14, bold=True, color=color)
    for j, item in enumerate(items):
        txt(s, f"• {item}", 0.45 + i*3.27, 2.15 + j*0.72, 2.9, 0.65,
            size=12, color=WHITE if j == 0 else MUTED, wrap=True)
    # 사용 시점
    use_bg = s.shapes.add_shape(1, Inches(0.35 + i*3.27), Inches(5.3), Inches(3.0), Inches(1.4))
    use_bg.fill.solid(); use_bg.fill.fore_color.rgb = CARD2; use_bg.line.fill.background()
    txt(s, "📌 언제 쓰나", 0.5 + i*3.27, 5.38, 2.8, 0.38, size=11, bold=True, color=color)
    txt(s, use_case,     0.5 + i*3.27, 5.75, 2.8, 0.85, size=12, color=WHITE, align=PP_ALIGN.CENTER)

tip_box(s, 0.3, 7.02, 12.7,
        "HTML 보고서는 브라우저에서 열리며 인터넷 연결 없이도 동작합니다. 팀원에게 파일로 전달하면 동일하게 확인 가능합니다.")


# ════════════════════════════════════════════════════════════════
# Slide 9 — 실험 전 체크리스트
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s, ORANGE)
txt(s, "실험 전 최종 체크리스트", 0.4, 0.1, 12, 0.4, size=13, color=ORANGE, bold=True)
txt(s, "프라이머 합성 발주 전 반드시 확인할 것들", 0.4, 0.48, 12, 0.6, size=24, bold=True)

checklist = [
    (ACCENT,  "커버리지",    "선택 후보의 커버리지 ≥ 80% 확인",
               "분자진단 제품은 변이 대응률이 핵심 — 낮으면 실패"),
    (PINK,    "Tm 균형",     "포워드·리버스 Tm 차이 ≤ 2°C 확인",
               "차이가 크면 한 프라이머가 먼저 탈락해 증폭 불안정"),
    (YELLOW,  "GC%",         "두 프라이머 모두 40–60% 범위 확인",
               "GC 극단값은 비특이 결합 및 이합체 형성 위험"),
    (BLUE,    "AI 점수",     "최종 선택 후보 AI 점수 ≥ 60점 확인",
               "60점 미만이면 근거 세부 확인 후 재설계 검토"),
    (PURPLE,  "이합체",      "3'-end 이합체 Δ G > -6 kcal/mol 확인",
               "이합체가 안정적이면 실제 PCR에서 artifact 발생"),
    (ORANGE,  "제품 크기",   "qPCR: 80–200 bp, PCR: 목적에 맞는 크기 확인",
               "크기가 너무 크면 qPCR 효율 급락"),
    (ACCENT2, "서열 중복",   "동일 서열 후보 여러 개 발주 않도록 확인",
               "이웃한 후보는 서열이 거의 동일할 수 있음"),
    (MUTED,   "보고서 저장", "PDF 또는 HTML 보고서 ELN에 첨부",
               "실험 재현성 및 규제 문서화 요건 대응"),
]

for i, (color, title, check, reason) in enumerate(checklist):
    row = i % 4
    col = i // 4
    bx = box(s, 0.3 + col*6.6, 1.35 + row*1.42, 6.3, 1.3, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.3 + col*6.6), Inches(1.35 + row*1.42),
                             Inches(0.06), Inches(1.3))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    # 체크박스 모양
    chk = s.shapes.add_shape(1, Inches(0.5 + col*6.6), Inches(1.5 + row*1.42),
                              Inches(0.3), Inches(0.3))
    chk.fill.solid(); chk.fill.fore_color.rgb = CARD2
    chk.line.color.rgb = color
    txt(s, title, 0.95 + col*6.6, 1.47 + row*1.42, 3.0, 0.4, size=14, bold=True, color=color)
    txt(s, check, 0.5  + col*6.6, 1.9  + row*1.42, 5.7, 0.38, size=12, color=WHITE)
    txt(s, f"→ {reason}", 0.5 + col*6.6, 2.27 + row*1.42, 5.7, 0.3, size=10.5, color=MUTED, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 10 — 자주 묻는 질문 (FAQ)
# ════════════════════════════════════════════════════════════════
s = slide(); bg(s)
top_bar(s, MUTED)
txt(s, "FAQ",  0.4, 0.1, 5, 0.4, size=13, color=MUTED, bold=True)
txt(s, "시약 연구원이 자주 묻는 질문", 0.4, 0.48, 12, 0.6, size=24, bold=True)

faqs = [
    ("Q. FASTA 서열이 적으면 결과를 믿을 수 있나요?",
     "서열 10개 미만이면 커버리지 통계가 부정확합니다. 가능하면 20개 이상, 대표 변이주를 포함한 서열을 사용하세요.",
     BLUE),
    ("Q. AI 점수 70점인 후보를 발주해도 될까요?",
     "AI 점수 단독 판단은 금물입니다. 커버리지·Tm·최종 점수를 함께 확인하세요. 70점이라도 커버리지 95%·Tm 균형이면 우수한 후보입니다.",
     ACCENT),
    ("Q. 커버리지가 60%밖에 안 나옵니다.",
     "① Max Mismatches를 3으로 완화  ② 보존 영역 임계값 낮추기  ③ 타겟 서열이 너무 다양하면 multiplex 설계로 전환을 고려하세요.",
     YELLOW),
    ("Q. Tm이 65°C를 넘는 후보만 나옵니다.",
     "Advanced Options에서 Primer Tm 최대값을 조정하거나, 보존 영역을 GC 함량이 낮은 부위로 좁혀보세요.",
     ORANGE),
    ("Q. Multiplex용 프로브 서열은 어디서 확인하나요?",
     "결과 목록의 'Probe' 컬럼 또는 HTML 보고서 서열 섹션에서 확인합니다. 어세이 타입을 Multiplex qPCR로 선택해야 생성됩니다.",
     PURPLE),
    ("Q. 결과 보고서를 동료에게 공유하려면?",
     "보고서 화면에서 PDF 다운로드 버튼을 누르거나, HTML 파일을 그대로 이메일로 전달하면 수신자도 동일하게 확인 가능합니다.",
     PINK),
]

for i, (q, a, color) in enumerate(faqs):
    row = i % 3
    col = i // 3
    bx = box(s, 0.3 + col*6.6, 1.35 + row*1.92, 6.3, 1.8, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.3 + col*6.6), Inches(1.35 + row*1.92),
                             Inches(6.3), Inches(0.06))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s, q, 0.5 + col*6.6, 1.45 + row*1.92, 5.9, 0.52, size=12.5, bold=True, color=color, wrap=True)
    txt(s, a, 0.5 + col*6.6, 1.98 + row*1.92, 5.9, 1.0,  size=12,   color=MUTED, wrap=True)


# ── 저장 ──────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "OpenBioShield_연구원가이드.pptx")
prs.save(out)
print(f"저장 완료: {out}")
