"""
OpenBioShield Phase 1/2/3 — 시약 연구원 가이드 PPT
씨젠 e-Sign 디자인 시스템 기반
"""
import os, sys
sys.path.insert(0, r'C:\source\Bio-Visualizers')

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from ppt_theme import *

OUT = r'C:\source\Bio-Visualizers'

PHASE1_ACCENT = CYAN    # #06B6D4
PHASE2_ACCENT = WARN    # #D97706
PHASE3_ACCENT = SUCCESS # #16A34A


# ════════════════════════════════════════════════════════════════
#  PHASE 1  —  DNA 탐색 & 변이 분석
# ════════════════════════════════════════════════════════════════
def build_phase1():
    prs = new_prs()
    A = PHASE1_ACCENT

    # ── S1 타이틀 ──────────────────────────────────────────────
    make_title_slide(prs, "1",
        "Phase 1  |  시약 연구원 가이드",
        "DNA 탐색 & 변이 분석",
        "COVID-19 역학 대시보드  ·  유전체 브라우저  ·  프라이머 구조 분석",
        accent=A)

    # ── S2 개요 — 3열 다크카드 ─────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_DARK)
    make_content_header(s, "01", "Phase 1 개요", "3개 화면으로 구성된 탐색 워크플로우", accent=A)
    dark_card_3col(s, 0.35, 1.18, 4.09, 4.9, "01", "COVID 대시보드",
        ["역학 컨텍스트 파악", "국가·지역별 발생 현황", "CFR 트렌드 차트", "타겟 질환 선택"], A)
    dark_card_3col(s, 4.62, 1.18, 4.09, 4.9, "02", "DNA 서열 맵",
        ["NCBI 참조 서열 자동 로드", "변이 위치 오버레이", "관심 영역 줌인", "프라이머 결합 위치 확인"], A)
    dark_card_3col(s, 8.89, 1.18, 4.09, 4.9, "03", "프라이머 구조",
        ["2차 구조 아크 다이어그램", "GC% · Tm 수치 확인", "변이주별 미스매치 검출", "왓슨-크릭 결합 시각화"], A)
    rect(s, 0.35, 6.22, 12.63, 0.72, CARD_DARK)
    rect(s, 0.35, 6.22, 0.06,  0.72, A)
    textbox(s, "활용 순서", 0.5, 6.3, 1.5, 0.3, size=11, bold=True, color=A)
    textbox(s, "COVID 대시보드(역학 파악) → DNA 맵(보존 영역 확인) → 프라이머 구조(설계 적합성 평가) 순으로 진행하면 가장 효율적입니다.",
            2.1, 6.3, 10.7, 0.5, size=11, color=MUTED, wrap=True)

    # ── S3 COVID 대시보드 ──────────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "1", "COVID 대시보드", "역학 컨텍스트 파악 — 어떤 변이주를 타겟할지 결정하는 첫 단계", accent=A)

    # 왼쪽 단계 카드
    steps_l = [
        ("1", "세계 지도 확인", "국가별 확진자·사망자 색상 히트맵\n관심 국가 클릭 → 상세 수치 확인"),
        ("2", "대한민국 시도별 현황", "시도별 발생 현황 → 국내 진단 수요 예측에 활용"),
        ("3", "치명률(CFR) 트렌드", "CFR 높고 확산 중인 변이주 확인 → 개발 우선순위 결정"),
        ("4", "질환 선택", "SARS-CoV-2 외 HPV, STI 등 타 질환으로 전환 가능"),
    ]
    for i, (n, title, desc) in enumerate(steps_l):
        step_card(s, 0.35, 1.18+i*1.35, 6.2, 1.22, n, A, title, desc)

    # 오른쪽 활용 시나리오
    textbox(s, "실무 활용 시나리오", 6.85, 1.18, 6.1, 0.4,
            size=13, bold=True, color=TEXT_DARK)
    scenarios = [
        ("신규 진단 키트 타겟 선정", A,
         "CFR·확산 속도 데이터로 개발 필요성 높은 변이주 선택\n→ 해당 서열로 Phase 3 진행"),
        ("국내 특화 진단 제품", WARN,
         "시도별 지도에서 유행 지역 확인\n→ 지역 특화 변이주 타겟 설계"),
        ("개발 배경 근거 문서화", SUCCESS,
         "역학 데이터 스크린샷 → ELN·보고서에 첨부\n→ 허가 신청 배경 자료"),
    ]
    for i, (title, color, desc) in enumerate(scenarios):
        rect(s, 6.85, 1.68+i*1.65, 6.1, 1.52, WHITE, BORDER)
        rect(s, 6.85, 1.68+i*1.65, 0.06, 1.52, color)
        textbox(s, title, 7.05, 1.78+i*1.65, 5.8, 0.4, size=13, bold=True, color=color)
        textbox(s, desc,  7.05, 2.22+i*1.65, 5.8, 0.85, size=11, color=SLATE, wrap=True)
    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "대시보드는 실시간 데이터를 반영합니다. 개발 착수 전 최신 유행 상황을 확인해 타겟 변이주를 최신화하세요.", A)

    # ── S4 DNA 서열 맵 ─────────────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "2", "DNA 서열 맵 (유전체 브라우저)",
                        "변이 위치를 유전자 지도에서 직접 확인하고 프라이머 설계 타겟 구간을 선정합니다", accent=A)

    # 상단 4가지 기능
    funcs = [
        ("NCBI 자동 로드", "인터넷 연결 시\n최신 참조 서열 수신"),
        ("변이 위치 표시", "SNP·InDel 위치에\n색상 마커 자동 표시"),
        ("영역 줌인", "관심 유전자 영역\n드래그로 확대"),
        ("프라이머 위치", "기존 프라이머\n결합 구간 표시"),
    ]
    for i, (title, desc) in enumerate(funcs):
        rect(s, 0.35+i*3.27, 1.18, 3.1, 1.38, WHITE, BORDER)
        rect(s, 0.35+i*3.27, 1.18, 3.1, 0.04, A)
        textbox(s, title, 0.5+i*3.27, 1.28, 2.8, 0.42,
                size=13, bold=True, color=A)
        textbox(s, desc, 0.5+i*3.27, 1.72, 2.8, 0.75,
                size=11, color=SLATE, wrap=True)

    # 단계별 사용법
    textbox(s, "단계별 사용법", 0.35, 2.75, 12.63, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    dna_steps = [
        ("1", BLUE, "변이주 선택",
         "상단 드롭다운에서 분석할 변이주 선택 → 돌연변이 목록 자동 표시"),
        ("2", A,    "타겟 유전자 확인",
         "스파이크·RdRp·N 단백질 등 진단 타겟 영역 확인 — 변이 밀도 낮은 구간이 설계에 유리"),
        ("3", WARN, "보존 영역 메모",
         "변이 마커가 없는 구간 = 보존 영역 → Phase 3 FASTA 업로드 시 해당 영역 서열 준비"),
        ("4", SUCCESS, "프라이머 위치 검토",
         "기존 공개 프라이머와 신규 변이의 오버랩 확인 → 오버랩 시 재설계 필요"),
    ]
    for i, (n, color, title, desc) in enumerate(dna_steps):
        row, col = i % 2, i // 2
        step_card(s, 0.35+col*6.52, 3.22+row*1.48, 6.32, 1.35, n, color, title, desc)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "변이 마커가 밀집한 핫스팟은 피하고, 여러 변이주에서 공통으로 비어있는 구간을 프라이머 타겟으로 선택하세요.", A)

    # ── S5 프라이머 구조 분석 ──────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "3", "프라이머 구조 분석",
                        "기존 프라이머의 2차 구조·Tm·변이 영향을 실험 전에 사전 평가합니다", accent=A)

    # 왼쪽: 수치 판독 기준 테이블
    textbox(s, "수치 판독 기준", 0.35, 1.18, 5.9, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    headers_t = ["지표", "이상 범위", "판정"]
    col_widths = [1.5, 1.8, 2.35]
    col_xs = [0.35, 1.85, 3.65]
    for ci, (h, cw, cx) in enumerate(zip(headers_t, col_widths, col_xs)):
        rect(s, cx, 1.62, cw, 0.45, BLUE)
        textbox(s, h, cx+0.05, 1.65, cw-0.1, 0.38,
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    metrics = [
        ("GC%",   "40–60%",          "이상 범위 권장"),
        ("Tm",    "58–65°C",         "qPCR 권장"),
        ("MFE",   "> -6 kcal/mol",   "구조 형성 위험 없음"),
        ("미스매치", "0–1개",          "2개 이상 → 재설계"),
    ]
    for ri, (metric, ideal, verdict) in enumerate(metrics):
        bg_c = GRAY_BG if ri % 2 == 0 else WHITE
        for ci, (val, cw, cx) in enumerate(zip([metric, ideal, verdict], col_widths, col_xs)):
            rect(s, cx, 2.12+ri*0.55, cw, 0.52, bg_c, BORDER)
            c = A if ci == 0 else (TEXT_DARK if ci == 1 else SLATE)
            b = (ci == 0)
            textbox(s, val, cx+0.07, 2.17+ri*0.55, cw-0.1, 0.42,
                    size=12, bold=b, color=c, align=PP_ALIGN.CENTER)

    # 왼쪽 하단: 화면 구성
    textbox(s, "화면 구성 요소", 0.35, 4.48, 5.9, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    viz_items = [
        ("아크 다이어그램",  "프라이머 2차 구조 — 아크 밀집 = 헤어핀 위험"),
        ("염기 조성 막대",   "A·T·G·C 위치별 비율 — GC 편중 구간 파악"),
        ("변이주 미스매치",  "선택 변이주의 미스매치 위치 빨간 마커 표시"),
    ]
    for i, (name, desc) in enumerate(viz_items):
        rect(s, 0.35, 4.93+i*0.62, 5.9, 0.55, WHITE, BORDER)
        rect(s, 0.35, 4.93+i*0.62, 0.06, 0.55, A)
        textbox(s, name, 0.55, 4.98+i*0.62, 2.0, 0.38, size=12, bold=True, color=A)
        textbox(s, desc, 2.6, 4.98+i*0.62, 3.5, 0.38, size=11, color=SLATE, wrap=True)

    # 오른쪽: 활용 케이스
    textbox(s, "실무 활용 케이스", 6.55, 1.18, 6.45, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    cases = [
        ("기존 키트 변이 영향 평가", A,
         "현재 시판 프라이머 서열 입력 → 신규 변이주 선택\n→ 미스매치 개수 확인 → 2개 이상이면 재설계 착수"),
        ("Phase 3 결과 교차 검증", BLUE,
         "Phase 3에서 설계된 신규 프라이머를 이 화면에서\n2차 구조·Tm 재확인 후 실험 진행 여부 최종 결정"),
        ("MFE 주의 기준", WARN,
         "MFE < -6 kcal/mol: 헤어핀 구조 가능성 높음\n→ 해당 프라이머는 Phase 3 재설계 권장"),
        ("3'-end 미스매치 판단", RGBColor(0xDC,0x26,0x26),
         "3' 말단 3bp 내 미스매치 존재 시\n→ PCR 실패 위험 매우 높음 → 즉각 재설계"),
    ]
    for i, (title, color, desc) in enumerate(cases):
        rect(s, 6.55, 1.65+i*1.42, 6.45, 1.3, WHITE, BORDER)
        rect(s, 6.55, 1.65+i*1.42, 0.06, 1.3, color)
        textbox(s, title, 6.75, 1.74+i*1.42, 6.1, 0.38, size=13, bold=True, color=color)
        textbox(s, desc,  6.75, 2.16+i*1.42, 6.1, 0.72, size=11, color=SLATE, wrap=True)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "MFE 값이 -6 kcal/mol보다 낮으면(더 안정적이면) 헤어핀 구조 가능성이 높습니다. 해당 프라이머는 Phase 3 재설계를 권장합니다.", A)

    # ── S6 변이주 선택 가이드 ─────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "★", "변이주 선택 가이드",
                        "세 화면에서 공통으로 사용하는 변이주 선택기 — 정확한 선택이 분석 품질을 결정합니다", accent=A)

    # 왼쪽: 변이주 목록
    textbox(s, "지원 변이주 (SARS-CoV-2 기준)", 0.35, 1.18, 6.0, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    variants = [
        ("Wild-type",    "우한 원형 서열",                   MUTED),
        ("Alpha",        "B.1.1.7  |  N501Y  |  영국발",     BLUE),
        ("Beta",         "B.1.351  |  E484K  |  남아공발",    A),
        ("Delta",        "B.1.617.2  |  L452R  |  인도발",    WARN),
        ("Omicron BA.1", "돌연변이 32개  |  기존 프라이머 영향 큼",
                                                              RGBColor(0xDC,0x26,0x26)),
        ("Omicron BA.5", "국내 주요 유행 아형",              RGBColor(0x79,0x16,0xFA)),
        ("XBB / JN.1",   "최신 재조합 변이주",               RGBColor(0x0E,0xA5,0xE9)),
    ]
    for i, (name, desc, color) in enumerate(variants):
        rect(s, 0.35, 1.65+i*0.68, 6.0, 0.62, WHITE, BORDER)
        rect(s, 0.35, 1.65+i*0.68, 0.06, 0.62, color)
        textbox(s, name, 0.55, 1.7+i*0.68, 1.8, 0.42, size=12, bold=True, color=color)
        textbox(s, desc, 2.45, 1.7+i*0.68, 3.75, 0.42, size=11, color=SLATE)

    # 오른쪽 상단: 지원 질환
    textbox(s, "지원 질환", 6.6, 1.18, 6.35, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    diseases = [
        ("SARS-CoV-2", RGBColor(0xDC,0x26,0x26), "COVID-19 진단"),
        ("HPV",        WARN,                       "자궁경부암 관련"),
        ("STI",        BLUE,                       "성매개 감염 패널"),
    ]
    for i, (d, c, sub) in enumerate(diseases):
        rect(s, 6.6, 1.65+i*0.82, 6.35, 0.72, WHITE, BORDER)
        rect(s, 6.6, 1.65+i*0.82, 0.06, 0.72, c)
        textbox(s, d,   6.8,  1.73+i*0.82, 1.8, 0.42, size=13, bold=True, color=c)
        textbox(s, sub, 8.65, 1.73+i*0.82, 4.2, 0.42, size=12, color=SLATE)

    # 오른쪽 하단: 실무 팁
    textbox(s, "변이주 선택 실무 팁", 6.6, 4.2, 6.35, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    tips = [
        (A,    "현재 국내 유행 아형부터 분석 시작"),
        (BLUE, "기존 키트 평가 시 → 출시 당시 + 최신 변이 두 가지 비교"),
        (WARN, "Omicron BA.1은 스파이크 프라이머에 가장 큰 영향 → 반드시 재검토"),
        (MUTED,"신규 변이 추가는 개발팀 요청으로 반영 가능"),
    ]
    for i, (color, tip) in enumerate(tips):
        rect(s, 6.6, 4.65+i*0.6, 6.35, 0.52, INFO_BG, BORDER)
        rect(s, 6.6, 4.65+i*0.6, 0.06, 0.52, color)
        textbox(s, tip, 6.8, 4.7+i*0.6, 6.05, 0.42, size=12, color=TEXT_DARK, wrap=True)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "Omicron BA.1은 스파이크 영역에 돌연변이가 32개입니다. 스파이크 타겟 프라이머를 보유 중이라면 반드시 재검토하세요.", A)

    # ── S7 Phase 1→2→3 연결 ───────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_DARK)
    make_content_header(s, "→", "Phase 1 활용 결론",
                        "Phase 1에서 얻은 정보를 Phase 2, 3에 어떻게 연결하는가", accent=A)
    cols = [
        ("Phase 1 산출물", A, [
            "타겟 변이주 결정 (역학 데이터 기반)",
            "프라이머 결합 후보 구간 (DNA 맵 보존 영역)",
            "기존 프라이머 재설계 필요 여부 (미스매치 확인)",
            "개발 배경 근거 자료 (역학 스크린샷)",
        ]),
        ("→  Phase 2 전달", WARN, [
            "기존 프라이머 서열 → EP05/EP09 정밀도 검증",
            "변이 위치 정보 → PFPS 위험도 계산 입력값",
            "미스매치 위치 → 3'-end 근접 여부 확인",
            "역학 데이터 → 검증 우선순위 결정",
        ]),
        ("→  Phase 3 전달", SUCCESS, [
            "보존 영역 좌표 → FASTA 구성 기준",
            "타겟 변이주 서열 → FASTA 포함 목록",
            "어세이 타입 결정 → qPCR/Multiplex 선택",
            "Tm 기준 파악 → Advanced Options 설정",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        textbox(s, title, 0.35+i*4.35, 1.3, 4.0, 0.42,
                size=14, bold=True, color=color)
        for j, item in enumerate(items):
            rect(s, 0.35+i*4.35, 1.82+j*1.1, 4.09, 0.95, CARD_DARK)
            rect(s, 0.35+i*4.35, 1.82+j*1.1, 0.06, 0.95, color)
            textbox(s, item, 0.55+i*4.35, 1.9+j*1.1, 3.75, 0.78,
                    size=12, color=WHITE, wrap=True)

    prs.save(os.path.join(OUT, 'OpenBioShield_Phase1_연구원가이드_v2.pptx'))
    print("Phase 1 저장 완료")


# ════════════════════════════════════════════════════════════════
#  PHASE 2  —  통계 검증 & 위험도 분석
# ════════════════════════════════════════════════════════════════
def build_phase2():
    prs = new_prs()
    A = PHASE2_ACCENT   # amber #D97706

    # ── S1 타이틀 ──────────────────────────────────────────────
    make_title_slide(prs, "2",
        "Phase 2  |  시약 연구원 가이드",
        "통계 검증 & 위험도 분석",
        "CLSI EP05/EP09  ·  PFPS 위험도 엔진  ·  데이터 수집",
        accent=A)

    # ── S2 개요 3열 다크카드 ───────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_DARK)
    make_content_header(s, "01", "Phase 2 개요", "3개 화면으로 구성된 검증 & 위험도 워크플로우", accent=A)
    dark_card_3col(s, 0.35, 1.18, 4.09, 4.9, "01", "CLSI 통계 검증",
        ["EP05 반복성·재현성 분산분석", "EP09 방법비교 Deming 회귀", "AI 보고서 자동 생성", "결과 DB 자동 저장"], A)
    dark_card_3col(s, 4.62, 1.18, 4.09, 4.9, "02", "PFPS 위험도 엔진",
        ["변이 위치 거리 기반 위험 점수", "3단계 위험 등급 (HIGH/MEDIUM/LOW)", "GPT-4o 서술형 설명", "왓슨-크릭 결합 시뮬레이션"], A)
    dark_card_3col(s, 8.89, 1.18, 4.09, 4.9, "03", "데이터 수집 현황",
        ["수집된 검증 레코드 통계", "질환·가이드라인별 분포", "Fine-tuning 준비도 지표", "연구원 피드백 집계"], A)
    rect(s, 0.35, 6.22, 12.63, 0.72, CARD_DARK)
    rect(s, 0.35, 6.22, 0.06,  0.72, A)
    textbox(s, "규제 대응", 0.5, 6.3, 1.4, 0.3, size=11, bold=True, color=A)
    textbox(s, "EP05·EP09 결과와 AI 보고서는 식약처(MFDS)·CE-IVD(IVDR)·FDA 510(k) 허가 자료로 직접 활용 가능합니다.",
            2.0, 6.3, 11.0, 0.5, size=11, color=MUTED, wrap=True)

    # ── S3 EP05 단계별 가이드 ──────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "1", "CLSI EP05 — 정밀도 검증",
                        "Excel 파일 업로드 한 번으로 반복성·재현성 ANOVA, CV% 자동 계산", accent=A)

    # 왼쪽: 단계 카드
    steps_ep = [
        ("1", "Excel 파일 업로드",
         "CLSI EP05 형식 데이터 파일 드래그&드롭\n(샘플 × 반복 × 재현 측정값)"),
        ("2", "컬럼 매핑 확인",
         "AI 자동 감지 컬럼 할당 검토\n잘못된 경우 드롭다운으로 직접 수정"),
        ("3", "EP05 분석 실행",
         "클릭 한 번 → ANOVA, CV%, 그랜드 평균 자동 계산\n완료까지 약 5–10초"),
        ("4", "AI 보고서 & DB 저장",
         "한국어·영어 전환 가능한 AI 마크다운 보고서\n분석 결과가 Supabase DB에 자동 저장"),
    ]
    for i, (n, title, desc) in enumerate(steps_ep):
        step_card(s, 0.35, 1.18+i*1.38, 6.1, 1.25, n, A, title, desc)

    # 오른쪽: 수치 해석 테이블 + 판정
    textbox(s, "EP05 결과 수치 해석", 6.65, 1.18, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    ep05_cols = ["지표", "허용 기준", "초과 시 조치"]
    col_w2 = [1.5, 1.7, 2.85]
    col_x2 = [6.65, 8.15, 9.85]
    for ci, (h, cw, cx) in enumerate(zip(ep05_cols, col_w2, col_x2)):
        rect(s, cx, 1.62, cw, 0.45, A)
        textbox(s, h, cx+0.05, 1.65, cw-0.1, 0.38,
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ep05_rows = [
        ("반복성 CV",   "< 5%",   "측정 조건 재검토 (온도, 균질도)"),
        ("재현성 CV",   "< 10%",  "다일·다기기 표준화 후 재분석"),
        ("ANOVA p값",  "< 0.05", "측정 그룹 간 조건 차이 확인"),
        ("그랜드 평균", "기준값 ±5%", "시약 농도·보정 재점검"),
    ]
    for ri, row in enumerate(ep05_rows):
        bg_c = GRAY_BG if ri % 2 == 0 else WHITE
        for ci, (val, cw, cx) in enumerate(zip(row, col_w2, col_x2)):
            rect(s, cx, 2.12+ri*0.55, cw, 0.52, bg_c, BORDER)
            c = A if ci == 0 else (TEXT_DARK if ci == 1 else SLATE)
            textbox(s, val, cx+0.07, 2.17+ri*0.55, cw-0.1, 0.42,
                    size=11, color=c, align=PP_ALIGN.CENTER)

    textbox(s, "합격 / 불합격 판정", 6.65, 4.48, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    judgments = [
        (SUCCESS_BG, SUCCESS, "합격",     "반복성 CV < 5%  AND  재현성 CV < 10%  →  다음 단계 진행"),
        (WARN_BG,    WARN,    "조건부",   "CV 5–15%  →  추가 검증 후 사용 여부 결정"),
        (RGBColor(0xFF,0xF1,0xF2), RGBColor(0xDC,0x26,0x26), "불합격",
         "CV > 15% 또는 p > 0.05  →  프로토콜 재검토"),
    ]
    for i, (bg_c, color, verdict, cond) in enumerate(judgments):
        rect(s, 6.65, 4.93+i*0.65, 6.3, 0.58, bg_c, BORDER)
        rect(s, 6.65, 4.93+i*0.65, 0.06, 0.58, color)
        risk_badge_inline(s, 6.75, 4.98+i*0.65, verdict, color)
        textbox(s, cond, 8.1, 4.98+i*0.65, 4.7, 0.42, size=11, color=SLATE, wrap=True)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "재현성 CV > 10% 이면 PFPS 엔진이 HIGH로 자동 에스컬레이션합니다. 반드시 위험도 분석과 연계해 확인하세요.", A)

    # ── S4 EP09 방법비교 ───────────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "2", "CLSI EP09 — 방법비교 검증",
                        "신규 진단법 vs 기준법 동등성 확인 — Deming 회귀·Bland-Altman·Pearson r", accent=A)

    # 왼쪽: 3가지 분석 항목
    items_ep09 = [
        ("Deming 회귀",        A,
         "기울기 1.0 ± 0.1,  절편 0 ± 10%\n두 방법 측정값 스케일이 동일한지 확인\n기울기가 1에서 멀수록 비례 오차 존재"),
        ("Bland-Altman 분석",  BLUE,
         "평균 차이(Bias) ± 5% 이내\n두 방법 간 체계적 차이를 산포도로 시각화\n95% 일치구간 허용 범위 내 확인"),
        ("Pearson 상관계수",   SUCCESS,
         "r ≥ 0.95 합격 기준\n두 방법 간 선형 관련성 강도\nr < 0.90 이면 방법 재검토 필요"),
    ]
    for i, (title, color, desc) in enumerate(items_ep09):
        rect(s, 0.35, 1.18+i*1.75, 6.1, 1.62, WHITE, BORDER)
        rect(s, 0.35, 1.18+i*1.75, 0.06, 1.62, color)
        textbox(s, title, 0.55, 1.27+i*1.75, 5.7, 0.42, size=14, bold=True, color=color)
        textbox(s, desc,  0.55, 1.73+i*1.75, 5.7, 0.95, size=12, color=SLATE, wrap=True)

    # 오른쪽: 활용 시나리오
    textbox(s, "실무 활용 시나리오", 6.65, 1.18, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    ep09_cases = [
        ("신규 PCR 키트 허가 신청", A,
         "기준법과 신규 키트 동시 측정 → EP09 업로드\n→ r≥0.95, Bias±5% 확인 → 허가 서류 AI 보고서 출력"),
        ("장비 교체 동등성 확인", BLUE,
         "구형 장비 vs 신형 장비 동일 샘플 측정\n→ Deming 기울기 1.0±0.1 → 교체 승인"),
        ("POC vs Lab 플랫폼 비교", SUCCESS,
         "현장 검사(POC) vs 검사실 표준 장비\n→ Bland-Altman Bias 허용 범위 확인 → POC 도입 승인"),
        ("시약 로트 변경 검증", WARN,
         "기존 로트 vs 신규 로트 비교\n→ 기울기·절편 범위 내 → 품질관리 문서 자동 생성"),
    ]
    for i, (title, color, desc) in enumerate(ep09_cases):
        rect(s, 6.65, 1.65+i*1.28, 6.3, 1.15, WHITE, BORDER)
        rect(s, 6.65, 1.65+i*1.28, 0.06, 1.15, color)
        textbox(s, title, 6.85, 1.73+i*1.28, 5.9, 0.38, size=13, bold=True, color=color)
        textbox(s, desc,  6.85, 2.15+i*1.28, 5.9, 0.55, size=11, color=SLATE, wrap=True)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "AI 보고서는 한국어·영어 전환 가능합니다. 허가 신청 시 영문 버전을 출력하면 해외 규제 제출 자료로 바로 활용할 수 있습니다.", A)

    # ── S5 PFPS 위험도 엔진 ────────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "3", "PFPS 위험도 엔진",
                        "프라이머-변이 충돌 위험을 위치 거리 기반 점수로 수치화합니다", accent=A)

    # 왼쪽: PFPS 개념 + 위험 구간
    info_box(s, 0.35, 1.18, 6.1, 1.38,
             "PFPS (Primer Fitness Position Score)",
             "변이 위치가 프라이머 결합 구간과 얼마나 가까운지를 거리 가중치로 계산한 위험 점수입니다.\n3'-end(3 프라임 말단)에 가까울수록 PCR 실패 위험이 기하급수적으로 높아집니다.", A)

    textbox(s, "위치별 위험 등급", 0.35, 2.72, 6.1, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    risk_zones = [
        ("0–3 bp",   "3'-end 내부",    "HIGH",   RGBColor(0xDC,0x26,0x26),
         "증폭 실패 확률 매우 높음 — 즉각 재설계"),
        ("4–15 bp",  "결합 중간 부위", "MEDIUM", WARN,
         "증폭 효율 저하 가능 — Gradient PCR 고려"),
        ("16+ bp",   "결합 외곽",      "LOW",    SUCCESS,
         "영향 미미 — 모니터링 수준 유지"),
    ]
    for i, (dist, zone, level, color, desc) in enumerate(risk_zones):
        rect(s, 0.35, 3.16+i*1.1, 6.1, 0.98, WHITE, BORDER)
        rect(s, 0.35, 3.16+i*1.1, 0.06, 0.98, color)
        textbox(s, dist,  0.55, 3.23+i*1.1, 1.1, 0.38, size=13, bold=True, color=color)
        textbox(s, zone,  1.7,  3.23+i*1.1, 1.6, 0.38, size=12, color=MUTED)
        risk_badge_inline(s, 3.4, 3.23+i*1.1, level, color)
        textbox(s, desc,  0.55, 3.65+i*1.1, 5.7, 0.38, size=11, color=SLATE, wrap=True)

    # 오른쪽: 결과 화면 구성 + 판단 기준
    textbox(s, "결과 화면 구성 요소", 6.65, 1.18, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    pfps_ui = [
        ("위험도 등급 카드",     A,    "HIGH/MEDIUM/LOW 색상 배지 + 위험 점수 수치"),
        ("결합 시뮬레이션",      BLUE, "프라이머↔타겟 서열 정렬 — 미스매치 위치 강조"),
        ("GPT-4o 설명",         RGBColor(0x79,0x16,0xFA), "위험 원인·개선 방향 서술형 해설"),
        ("변이 위치 목록",       MUTED,"각 미스매치의 bp 거리·영향도 상세 표"),
    ]
    for i, (name, color, desc) in enumerate(pfps_ui):
        rect(s, 6.65, 1.65+i*0.72, 6.3, 0.65, WHITE, BORDER)
        rect(s, 6.65, 1.65+i*0.72, 0.06, 0.65, color)
        textbox(s, name, 6.85, 1.72+i*0.72, 2.3, 0.38, size=12, bold=True, color=color)
        textbox(s, desc, 9.2,  1.72+i*0.72, 3.6, 0.38, size=11, color=SLATE)

    textbox(s, "PFPS 결과 → 실험 판단 기준", 6.65, 4.65, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    pfps_actions = [
        ("HIGH",   RGBColor(0xDC,0x26,0x26), "즉각 재설계 → Phase 3으로 신규 프라이머 설계"),
        ("MEDIUM", WARN,                      "실험 진행 + Gradient PCR 병행 → 최적화 검증"),
        ("LOW",    SUCCESS,                   "사용 적합 — 분기별 변이 업데이트 확인 권장"),
    ]
    for i, (level, color, action) in enumerate(pfps_actions):
        rect(s, 6.65, 5.1+i*0.65, 6.3, 0.58, WHITE, BORDER)
        rect(s, 6.65, 5.1+i*0.65, 0.06, 0.58, color)
        risk_badge_inline(s, 6.75, 5.16+i*0.65, level, color)
        textbox(s, action, 8.1, 5.16+i*0.65, 4.7, 0.42, size=12, color=SLATE, wrap=True)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "재현성 CV > 10% 이면 PFPS 엔진이 HIGH로 자동 에스컬레이션합니다. EP05 결과와 PFPS를 반드시 함께 확인하세요.", A)

    # ── S6 데이터 수집 & AI 개선 사이클 ───────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "4", "데이터 수집 현황",
                        "검증 데이터 축적 → AI 모델 자동 개선 — 연구원의 피드백이 핵심입니다", accent=A)

    # 왼쪽: 수집 화면 구성
    textbox(s, "데이터 수집 화면 구성 요소", 0.35, 1.18, 6.1, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    dc_items = [
        ("수집 레코드 수",   A,    "EP05·EP09 분석 완료 시 자동 저장된 총 레코드 수\n목표: 5,000건 이상 (Fine-tuning 기준)"),
        ("질환별 분포",      BLUE, "SARS-CoV-2·HPV·STI 각 질환별 데이터 비율 도넛 차트"),
        ("위험도 분포",      RGBColor(0xDC,0x26,0x26),
         "HIGH/MEDIUM/LOW 비율 — 치우침이 크면 데이터 균형화 필요"),
        ("피드백 정확도",    SUCCESS,
         "연구원 피드백 (accurate/partially/incorrect) 누적 — AI 보고서 품질 지표"),
        ("Fine-tuning 준비도", WARN,
         "5,000건 도달률 진행 바 — 목표 달성 시 XGBoost 재학습 알림"),
    ]
    for i, (title, color, desc) in enumerate(dc_items):
        rect(s, 0.35, 1.65+i*0.95, 6.1, 0.85, WHITE, BORDER)
        rect(s, 0.35, 1.65+i*0.95, 0.06, 0.85, color)
        textbox(s, title, 0.55, 1.72+i*0.95, 2.2, 0.38, size=13, bold=True, color=color)
        textbox(s, desc,  2.85, 1.72+i*0.95, 3.45, 0.65, size=11, color=SLATE, wrap=True)

    # 오른쪽: AI 개선 사이클
    textbox(s, "데이터 → AI 개선 사이클", 6.65, 1.18, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    cycle = [
        ("1", "EP05·EP09 분석 실행",    A,    "시약 연구원이 검증 데이터 업로드"),
        ("2", "결과 자동 저장",          BLUE, "분석 결과가 Supabase DB에 적재"),
        ("3", "피드백 입력",             WARN, "AI 보고서 정확도 3단계 평가"),
        ("4", "데이터 5,000건 도달",     RGBColor(0x79,0x16,0xFA), "Fine-tuning 준비 완료 알림"),
        ("5", "XGBoost 모델 재학습",     SUCCESS, "Phase 3 AI 점수 정확도 향상"),
        ("6", "더 정확한 자동 설계",     A,    "연구원 수작업 검토 시간 감소"),
    ]
    for i, (n, title, color, desc) in enumerate(cycle):
        rect(s, 6.65, 1.65+i*0.82, 6.3, 0.72, WHITE, BORDER)
        rect(s, 6.65, 1.65+i*0.82, 0.55, 0.72, color)
        textbox(s, n, 6.65, 1.65+i*0.82, 0.55, 0.72,
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, title, 7.3, 1.72+i*0.82, 2.5, 0.38, size=13, bold=True, color=color)
        textbox(s, desc,  9.85, 1.72+i*0.82, 2.95, 0.38, size=11, color=SLATE)

    success_box(s, 0.35, 6.88, 12.63, 0.5,
                "피드백 버튼(accurate·partially·incorrect)을 반드시 클릭해주세요. 이 피드백이 AI 모델 개선의 핵심 데이터입니다.")

    # ── S7 체크리스트 ─────────────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "✓", "Phase 2 완료 체크리스트",
                        "Phase 3 진행 전 반드시 완료해야 할 항목", accent=A)

    checks = [
        (A,    "EP05 완료",       "반복성 CV < 5%,  재현성 CV < 10%",
         "미달 시: 측정 조건 재검토 후 재분석"),
        (BLUE, "EP09 완료",       "r ≥ 0.95,  Deming 기울기 1.0±0.1,  Bias < 5%",
         "미달 시: 측정 범위·기준법 재검토"),
        (RGBColor(0xDC,0x26,0x26), "PFPS 확인", "HIGH 위험 프라이머 없음",
         "HIGH 존재 시: Phase 3으로 재설계 진행"),
        (WARN, "AI 보고서 피드백", "보고서 정확도 피드백 (3단계) 입력 완료",
         "피드백 없으면 DB 카운트에서 제외"),
        (SUCCESS, "DB 저장 확인", "데이터 수집 화면에서 오늘 레코드 수 증가 확인",
         "저장 실패 시: 네트워크 확인 또는 개발팀 문의"),
        (RGBColor(0x79,0x16,0xFA), "문서 저장", "규제 제출용 PDF·AI 보고서 ELN 첨부",
         "허가 신청 전 보고서 버전 관리 필수"),
    ]
    for i, (color, title, check, action) in enumerate(checks):
        row, col = i % 3, i // 3
        x, y = 0.35+col*6.55, 1.18+row*1.72
        rect(s, x, y, 6.3, 1.58, WHITE, BORDER)
        rect(s, x, y, 0.06, 1.58, color)
        rect(s, x+0.18, y+0.18, 0.42, 0.42, INFO_BG, BORDER)  # 체크박스
        textbox(s, title, x+0.75, y+0.18, 5.3, 0.42, size=13, bold=True, color=color)
        textbox(s, check, x+0.18, y+0.72, 5.95, 0.38, size=12, color=TEXT_DARK, wrap=True)
        textbox(s, f"→ {action}", x+0.18, y+1.12, 5.95, 0.35,
                size=11, color=SLATE, wrap=True, italic=True)

    prs.save(os.path.join(OUT, 'OpenBioShield_Phase2_연구원가이드_v2.pptx'))
    print("Phase 2 저장 완료")


# ════════════════════════════════════════════════════════════════
#  PHASE 3  —  자동 어세이 설계 파이프라인
# ════════════════════════════════════════════════════════════════
def build_phase3():
    prs = new_prs()
    A = PHASE3_ACCENT   # green #16A34A

    # ── S1 타이틀 ──────────────────────────────────────────────
    make_title_slide(prs, "3",
        "Phase 3  |  시약 연구원 가이드",
        "AI 기반 어세이 자동 설계",
        "9단계 파이프라인  ·  AI 효율 점수  ·  HTML/PDF 보고서",
        accent=A)

    # ── S2 9단계 파이프라인 개요 ───────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_DARK)
    make_content_header(s, "01", "9단계 자동화 파이프라인",
                        "FASTA 업로드부터 HTML/PDF 보고서까지 — 분 단위 완료", accent=A)

    pipeline = [
        ("1", "MSA\n정렬",      A,     "MAFFT"),
        ("2", "보존\n영역",     BLUE,  "Shannon"),
        ("3", "후보\n생성",     RGBColor(0x79,0x16,0xFA), "Primer3"),
        ("4", "특이성\n필터",   WARN,  "Off-target"),
        ("5", "커버리지",       RGBColor(0x0E,0xA5,0xE9), "변이 대응"),
        ("6", "열역학",         RGBColor(0xF4,0x72,0x7E), "Tm·GC·이합체"),
        ("7", "AI 점수",        A,     "효율 예측"),
        ("8", "순위화",         BLUE,  "가중 점수"),
        ("9", "보고서",         CYAN,  "HTML·PDF"),
    ]
    for i, (n, label, color, sub) in enumerate(pipeline):
        x = 0.35 + i*1.44
        rect(s, x, 1.2, 1.28, 2.5, CARD_DARK)
        rect(s, x, 1.2, 1.28, 0.06, color)
        circ = s.shapes.add_shape(9,
            __import__('pptx').util.Inches(x+0.43), __import__('pptx').util.Inches(1.3),
            __import__('pptx').util.Inches(0.42), __import__('pptx').util.Inches(0.42))
        circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
        textbox(s, n,     x+0.43, 1.3,  0.42, 0.42, size=14, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, label, x+0.04, 1.82, 1.2,  0.65, size=11, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, sub,   x+0.04, 2.52, 1.2,  0.55, size=9,
                color=MUTED, align=PP_ALIGN.CENTER)
        if i < 8:
            textbox(s, "→", x+1.28, 2.2, 0.18, 0.4, size=12,
                    color=MUTED, align=PP_ALIGN.CENTER)

    # 연구원 핵심 확인 포인트
    rect(s, 0.35, 3.88, 12.63, 2.92, CARD_DARK)
    rect(s, 0.35, 3.88, 0.06,  2.92, A)
    textbox(s, "시약 연구원이 집중할 단계", 0.55, 4.0, 12.0, 0.42,
            size=14, bold=True, color=RGBColor(0x6E,0xE7,0xB7))
    focus = [
        (A,    "Step 1",  "FASTA 파일 품질 — 서열 수가 많을수록 커버리지 정확도 향상 (20개 이상 권장)"),
        (RGBColor(0x0E,0xA5,0xE9), "Step 5",
         "커버리지 점수 — 80% 이상 후보만 실험 대상으로 선정 권장"),
        (RGBColor(0xF4,0x72,0x7E), "Step 6",
         "Tm 균형 — 포워드/리버스 Tm 차이 ±2°C 이내 후보 우선 선택"),
        (A,    "Step 7",  "AI 점수 근거 — 점수 낮은 이유를 보고서에서 클릭해 원인 파악"),
    ]
    for i, (color, step, desc) in enumerate(focus):
        circ2 = s.shapes.add_shape(9,
            __import__('pptx').util.Inches(0.55), __import__('pptx').util.Inches(4.55+i*0.55),
            __import__('pptx').util.Inches(0.35), __import__('pptx').util.Inches(0.35))
        circ2.fill.solid(); circ2.fill.fore_color.rgb = color; circ2.line.fill.background()
        textbox(s, step,  0.56, 4.55+i*0.55, 1.4, 0.35, size=11, bold=True, color=color)
        textbox(s, desc,  1.98, 4.55+i*0.55, 10.8, 0.35, size=12, color=WHITE)

    # ── S3 FASTA 업로드 & 설정 ─────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "1", "FASTA 업로드 & 파라미터 설정",
                        "좋은 입력이 좋은 결과를 만듭니다", accent=A)

    # 왼쪽: FASTA 파일 요건
    textbox(s, "FASTA 파일 요건", 0.35, 1.18, 6.1, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    fasta_reqs = [
        (A,    "권장 서열 수",   "20개 이상 — 커버리지 통계 신뢰도 확보"),
        (A,    "서열 길이",      "분석 타겟 유전자 전체 포함 (트리밍 불필요)"),
        (A,    "헤더 형식",      ">샘플명  또는  >accession_번호"),
        (WARN, "주의",          "중복 서열 포함 시 커버리지 과대 계산 가능"),
        (WARN, "주의",          "갭(-) 문자는 자동 제거 처리됨"),
    ]
    for i, (color, key, val) in enumerate(fasta_reqs):
        rect(s, 0.35, 1.65+i*0.78, 6.1, 0.7, WHITE, BORDER)
        rect(s, 0.35, 1.65+i*0.78, 0.06, 0.7, color)
        textbox(s, key, 0.55, 1.72+i*0.78, 1.7, 0.38, size=12, bold=True, color=color)
        textbox(s, val, 2.35, 1.72+i*0.78, 3.95, 0.38, size=12, color=SLATE)

    # 오른쪽: 어세이 타입 선택
    textbox(s, "어세이 타입 선택 가이드", 6.65, 1.18, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    assay_types = [
        ("qPCR",         A,    "정량 분석 · 형광 탐침 불필요\n제품 크기 80–200 bp 자동 최적화"),
        ("Multiplex qPCR", BLUE, "TaqMan 프로브 자동 설계 포함\n여러 타겟 동시 검출 시 선택"),
        ("Standard PCR", WARN, "겔 전기영동 확인용\n제품 크기 제한 완화"),
    ]
    for i, (atype, color, desc) in enumerate(assay_types):
        rect(s, 6.65, 1.65+i*1.62, 6.3, 1.48, WHITE, BORDER)
        rect(s, 6.65, 1.65+i*1.62, 0.06, 1.48, color)
        textbox(s, atype, 6.85, 1.75+i*1.62, 5.9, 0.45, size=14, bold=True, color=color)
        textbox(s, desc,  6.85, 2.23+i*1.62, 5.9, 0.75, size=12, color=SLATE, wrap=True)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "동일 타겟의 여러 변이주(variants) FASTA를 함께 포함하면 커버리지 분석의 신뢰도가 크게 올라갑니다.", A)

    # ── S4 결과 목록 해석 ──────────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "2", "후보 프라이머 결과 목록",
                        "각 컬럼이 무엇을 의미하는지 이해해야 올바른 후보를 선택할 수 있습니다", accent=A)

    # 가상 테이블
    headers_tbl = ["순위", "Forward 5'→3'", "Reverse 5'→3'", "Tm Fwd", "Tm Rev", "GC%", "커버리지", "AI점수", "최종점수"]
    col_ws = [0.5, 2.6, 2.6, 0.7, 0.7, 0.6, 0.85, 0.75, 0.82]
    col_xs = [0.35]
    for cw in col_ws[:-1]: col_xs.append(col_xs[-1]+cw+0.03)
    for ci, (h, cw, cx) in enumerate(zip(headers_tbl, col_ws, col_xs)):
        rect(s, cx, 1.18, cw, 0.42, A)
        textbox(s, h, cx+0.04, 1.21, cw-0.06, 0.36,
                size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=False)
    rows_data = [
        ["1","ATGGCTAGCTAGC…","CGATCGATCGAT…","62.3°","61.8°","52%","94.5%","78.4","88.2"],
        ["2","GCTAGCTAGCTAG…","TAGCTAGCTAGC…","60.1°","61.2°","48%","87.3%","71.2","82.5"],
        ["3","CGATCGATCGATC…","ATCGATCGATCG…","63.5°","60.2°","55%","79.1%","65.8","75.3"],
    ]
    row_bgs = [SUCCESS_BG, WHITE, GRAY_BG]
    for ri, row in enumerate(rows_data):
        for ci, (cell, cw, cx) in enumerate(zip(row, col_ws, col_xs)):
            rect(s, cx, 1.65+ri*0.52, cw, 0.48, row_bgs[ri], BORDER)
            c = A if (ri==0 and ci in [6,7,8]) else (TEXT_DARK if ri==0 else MUTED)
            textbox(s, cell, cx+0.04, 1.68+ri*0.52, cw-0.06, 0.4,
                    size=9.5, color=c, align=PP_ALIGN.CENTER, wrap=False)

    # 컬럼 설명 5개
    col_guides = [
        ("Tm (융해온도)",   A,    "Fwd·Rev 차이 ±2°C 이내\n차이 클수록 PCR 최적화 어려움"),
        ("GC%",             BLUE, "40–60% 범위 이상\n극단값은 비특이 결합 위험"),
        ("커버리지",        RGBColor(0x0E,0xA5,0xE9), "입력 서열 중 증폭 가능 비율\n80% 이상 후보 우선 선택"),
        ("AI 점수",         RGBColor(0x79,0x16,0xFA), "설계 품질 예측 (0–100)\n낮으면 근거 클릭해 원인 확인"),
        ("최종 점수",       WARN, "커버리지·Tm·AI·특이성 가중합\n실험 우선순위 결정 기준"),
    ]
    for i, (title, color, desc) in enumerate(col_guides):
        rect(s, 0.35+i*2.6, 3.7, 2.45, 2.45, WHITE, BORDER)
        rect(s, 0.35+i*2.6, 3.7, 2.45, 0.06, color)
        textbox(s, title, 0.5+i*2.6, 3.82, 2.2, 0.42, size=13, bold=True, color=color)
        textbox(s, desc,  0.5+i*2.6, 4.3,  2.2, 1.65, size=11, color=SLATE, wrap=True)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "커버리지 80% 미만이면 Advanced Options → Max Mismatches를 3으로 높이거나 보존 영역 임계값을 낮춰보세요.", A)

    # ── S5 열역학 + AI 점수 (2 in 1) ──────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "3", "열역학 분석 & AI 점수",
                        "실험 전 검증의 핵심 — Tm·GC·이합체 판독 기준과 AI 점수 근거 활용법", accent=A)

    # 왼쪽: 열역학
    textbox(s, "열역학 분석 판독 기준", 0.35, 1.18, 6.1, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    thermo = [
        ("Tm (융해온도)",    RGBColor(0xF4,0x72,0x7E),
         "qPCR 권장: 58–65°C\nFwd·Rev 차이 ≤2°C\nMultiplex: 프로브 Tm = 프라이머 +5~10°C"),
        ("GC% 함량",         BLUE,
         "이상: 40–60%\n3'-end 마지막 5nt: GC 2–3개 권장\nGC 클램프 >80% = 비특이 결합 위험"),
        ("이합체·헤어핀",    WARN,
         "ΔG > -6 kcal/mol 권장\n3'-end 이합체는 즉각 폐기\n헤어핀 Tm < 어닐링 온도"),
    ]
    for i, (title, color, desc) in enumerate(thermo):
        rect(s, 0.35, 1.65+i*1.62, 6.1, 1.48, WHITE, BORDER)
        rect(s, 0.35, 1.65+i*1.62, 0.06, 1.48, color)
        textbox(s, title, 0.55, 1.75+i*1.62, 5.7, 0.42, size=13, bold=True, color=color)
        textbox(s, desc,  0.55, 2.22+i*1.62, 5.7, 0.80, size=12, color=SLATE, wrap=True)

    # 오른쪽: AI 점수 구성
    textbox(s, "AI 점수 구성 (0–100점)", 6.65, 1.18, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    ai_comp = [
        ("기본 점수",        "+60점",   TEXT_DARK),
        ("GC% (양방향)",    "최대 -16점", RGBColor(0xDC,0x26,0x26)),
        ("Tm 균형",          "최대 -15점", RGBColor(0xDC,0x26,0x26)),
        ("3'-end GC 클램프", "최대 -13점", RGBColor(0xDC,0x26,0x26)),
        ("이중가닥 엔트로피","± 보정",   A),
        ("제품 크기",        "± 5점",    BLUE),
        ("보존 영역 품질",   "최대 -10점", WARN),
    ]
    cols_ai = [2.3, 1.5, 2.3]
    xs_ai   = [6.65, 8.95, 10.45]
    for ci, h in enumerate(["구성 요소", "점수 영향", "설명"]):
        rect(s, xs_ai[ci], 1.62, cols_ai[ci], 0.42, A)
        textbox(s, h, xs_ai[ci]+0.05, 1.65, cols_ai[ci]-0.1, 0.35,
                size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ai_notes = ["출발점","40–60% 벗어나면 감점","Tm차 1°C당 -3점",
                "너무 약/강하면 감점","복잡도 보너스","80–200bp +5점","낮을수록 유리"]
    for ri, ((title, delta, color), note) in enumerate(zip(ai_comp, ai_notes)):
        bg_c = GRAY_BG if ri % 2 == 0 else WHITE
        rect(s, xs_ai[0], 2.09+ri*0.55, cols_ai[0], 0.5, bg_c, BORDER)
        textbox(s, title, xs_ai[0]+0.07, 2.14+ri*0.55, cols_ai[0]-0.1, 0.4,
                size=11, bold=True, color=color)
        rect(s, xs_ai[1], 2.09+ri*0.55, cols_ai[1], 0.5, bg_c, BORDER)
        textbox(s, delta, xs_ai[1]+0.05, 2.14+ri*0.55, cols_ai[1]-0.1, 0.4,
                size=11, color=color, align=PP_ALIGN.CENTER)
        rect(s, xs_ai[2], 2.09+ri*0.55, cols_ai[2], 0.5, bg_c, BORDER)
        textbox(s, note, xs_ai[2]+0.05, 2.14+ri*0.55, cols_ai[2]-0.1, 0.4,
                size=10, color=SLATE)

    textbox(s, "AI 점수별 실험 판단", 6.65, 5.95, 6.3, 0.38,
            size=13, bold=True, color=TEXT_DARK)
    ai_j = [
        (SUCCESS, "80점↑", "바로 합성 주문"),
        (A,       "60–80점","보완 확인 후 진행"),
        (WARN,    "40–60점","설계 재검토 권장"),
        (RGBColor(0xDC,0x26,0x26),"40점↓","재설계 필수"),
    ]
    for i, (color, score, action) in enumerate(ai_j):
        rect(s, 6.65+i*1.68, 6.4, 1.55, 0.78, INFO_BG, BORDER)
        rect(s, 6.65+i*1.68, 6.4, 1.55, 0.06, color)
        textbox(s, score,  6.68+i*1.68, 6.5,  1.5, 0.32, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        textbox(s, action, 6.68+i*1.68, 6.84, 1.5, 0.3,  size=10, color=SLATE, align=PP_ALIGN.CENTER)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "Tm 차이가 3°C를 넘으면 Gradient PCR로 최적 어닐링 온도를 찾거나 다음 순위 후보로 교체하세요.", A)

    # ── S6 보고서 활용 ────────────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "4", "HTML 보고서 활용법",
                        "보고서의 각 섹션을 언제 어떻게 사용하는지 알면 실험 효율이 올라갑니다", accent=A)

    sections_rpt = [
        ("요약 테이블",    A,    "상위 후보 전체 수치 한눈에 비교\n최종 점수 기준 우선순위 확인",  "합성 주문 후보 결정"),
        ("서열 정보",      BLUE, "Forward·Reverse 전체 서열 복사\n프로브 서열 (Multiplex 시)",     "합성 발주 시 복사"),
        ("AI 점수 근거",   RGBColor(0x79,0x16,0xFA),
         "후보 클릭 → 점수 성분별 분해\n어느 특성이 점수 낮췄는지 파악",                          "점수 낮은 이유 파악"),
        ("PDF 다운로드",   WARN, "보고서 전체 PDF 저장\nELN 첨부·팀 공유용",                       "기록 보관·공유"),
    ]
    for i, (title, color, desc, usecase) in enumerate(sections_rpt):
        rect(s, 0.35+i*3.27, 1.18, 3.12, 5.38, WHITE, BORDER)
        rect(s, 0.35+i*3.27, 1.18, 3.12, 0.06, color)
        textbox(s, title, 0.5+i*3.27, 1.3,  2.85, 0.42, size=14, bold=True, color=color)
        textbox(s, desc,  0.5+i*3.27, 1.85, 2.85, 1.65, size=12, color=SLATE, wrap=True)
        rect(s, 0.42+i*3.27, 3.7,   2.98, 2.65, GRAY_BG)
        textbox(s, "사용 시점", 0.55+i*3.27, 3.82, 2.75, 0.38,
                size=11, bold=True, color=color)
        textbox(s, usecase, 0.55+i*3.27, 4.25, 2.75, 0.55,
                size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    info_box(s, 0.35, 6.88, 12.63, 0.5,
             "", "HTML 보고서는 브라우저만 있으면 인터넷 없이도 열립니다. 파일로 전달 시 수신자도 동일하게 확인 가능합니다.", A)

    # ── S7 실험 전 체크리스트 ─────────────────────────────────
    s = add_slide(prs); rect(s, 0, 0, 13.33, 7.5, BG_LIGHT)
    make_content_header(s, "✓", "실험 전 최종 체크리스트",
                        "프라이머 합성 발주 전 반드시 확인할 8가지", accent=A)

    chk_items = [
        (A,    "커버리지",    "선택 후보 커버리지 ≥ 80%",
         "낮으면 Max Mismatches 완화 또는 Multiplex 고려"),
        (RGBColor(0xF4,0x72,0x7E), "Tm 균형", "Fwd·Rev Tm 차이 ≤ 2°C",
         "차이 크면 Gradient PCR 최적화 후 결정"),
        (BLUE, "GC%",       "양 프라이머 모두 40–60%",
         "극단값은 비특이 결합·이합체 위험"),
        (RGBColor(0x79,0x16,0xFA), "AI 점수", "최종 선택 AI 점수 ≥ 60점",
         "낮으면 보고서 AI 근거 섹션에서 원인 확인"),
        (WARN, "이합체",    "3'-end 이합체 ΔG > -6 kcal/mol",
         "이합체 안정적이면 PCR artifact 발생 위험"),
        (RGBColor(0x0E,0xA5,0xE9), "제품 크기",
         "qPCR: 80–200 bp  /  PCR: 목적 크기",
         "크기 >250bp이면 qPCR 효율 급락"),
        (A,    "서열 중복", "유사 서열 후보 중복 발주 없는지 확인",
         "이웃한 후보는 서열이 거의 동일할 수 있음"),
        (MUTED,"보고서 저장", "PDF·HTML 보고서 ELN 첨부",
         "재현성·규제 문서화 요건 대응"),
    ]
    for i, (color, title, check, action) in enumerate(chk_items):
        row, col = i % 4, i // 4
        x, y = 0.35+col*6.55, 1.18+row*1.45
        rect(s, x, y, 6.3, 1.32, WHITE, BORDER)
        rect(s, x, y, 0.06, 1.32, color)
        rect(s, x+0.18, y+0.15, 0.38, 0.38, INFO_BG, BORDER)
        textbox(s, title, x+0.7, y+0.15, 5.4, 0.38, size=13, bold=True, color=color)
        textbox(s, check, x+0.18, y+0.62, 5.95, 0.32, size=12, color=TEXT_DARK)
        textbox(s, f"→ {action}", x+0.18, y+0.96, 5.95, 0.28,
                size=10.5, color=SLATE, italic=True, wrap=True)

    prs.save(os.path.join(OUT, 'OpenBioShield_Phase3_연구원가이드_v2.pptx'))
    print("Phase 3 저장 완료")


if __name__ == "__main__":
    build_phase1()
    build_phase2()
    build_phase3()
    print("모두 완료!")
