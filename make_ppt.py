"""
OpenBioShield 사내 성과 보고 PPT 생성 스크립트
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── 색상 팔레트 ──────────────────────────────────────────────
DARK_BG   = RGBColor(0x0F, 0x17, 0x2A)   # 네이비 배경
ACCENT    = RGBColor(0x10, 0xB9, 0x81)   # 에메랄드 그린
ACCENT2   = RGBColor(0x6E, 0xE7, 0xB7)   # 라이트 그린
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0x1E, 0x2D, 0x40)   # 카드 배경
MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # 서브텍스트
YELLOW    = RGBColor(0xFB, 0xBF, 0x24)   # 강조 숫자

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=DARK_BG):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def box(slide, x, y, w, h, fill=LIGHT_BG, radius=False):
    shape = slide.shapes.add_shape(
        17 if radius else 1,   # 17=rounded rect, 1=rect
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = 0.05
    return shape


def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def accent_bar(slide, y=0.18, h=0.06):
    shape = slide.shapes.add_shape(1, 0, Inches(y), Inches(1.2), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# 좌측 에메랄드 세로 바
bar = s.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

txt(s, "OpenBioShield", 0.35, 1.6, 9, 1.2, size=52, bold=True, color=WHITE)
txt(s, "AI 기반 분자진단 어세이 자동 설계 플랫폼", 0.35, 2.9, 11, 0.7,
    size=24, color=ACCENT2)
txt(s, "2026년 상반기 성과 보고", 0.35, 3.7, 8, 0.5, size=18, color=MUTED)

# 우측 장식 원
from pptx.util import Pt as _Pt
circ = s.shapes.add_shape(9, Inches(9.8), Inches(1.0), Inches(3.2), Inches(3.2))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x10, 0xB9, 0x81)
circ.line.fill.background()
# 투명도 효과를 위한 내부 원
circ2 = s.shapes.add_shape(9, Inches(10.2), Inches(1.4), Inches(2.4), Inches(2.4))
circ2.fill.solid(); circ2.fill.fore_color.rgb = DARK_BG
circ2.line.fill.background()
txt(s, "🧬", 10.5, 1.9, 2, 1.2, size=48, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# Slide 2 — 문제 & 솔루션
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "해결한 문제", 0.4, 0.1, 6, 0.5, size=14, color=ACCENT, bold=True)
txt(s, "기존 방식의 한계, 우리가 바꾼다", 0.4, 0.45, 10, 0.7, size=28, bold=True)

problems = [
    ("⏱", "설계 시간 과다", "전문가가 수작업으로 수일~수주 소요"),
    ("❌", "높은 실패율",   "실험 전 검증 어려워 반복 실패"),
    ("💸", "고비용",        "시약·인력·재실험 비용 누적"),
]
arrows = ["→", "→", "→"]
solutions = [
    ("✅", "완전 자동화",   "9단계 파이프라인, 분 단위 완료"),
    ("🎯", "사전 AI 검증",  "AI 효율 점수 + 커버리지 사전 필터"),
    ("☁", "클라우드 즉시", "Vercel 서버리스 — 인프라 불필요"),
]

for i, (ic, title, desc) in enumerate(problems):
    bx = box(s, 0.3, 1.4 + i*1.55, 4.0, 1.3, fill=RGBColor(0x1A, 0x1F, 0x35), radius=True)
    txt(s, ic,    0.45, 1.45 + i*1.55, 0.8, 0.8, size=28)
    txt(s, title, 1.2,  1.45 + i*1.55, 3.0, 0.45, size=16, bold=True)
    txt(s, desc,  1.2,  1.85 + i*1.55, 2.9, 0.45, size=12, color=MUTED)

for i, arrow in enumerate(arrows):
    txt(s, "➜", 4.5, 1.75 + i*1.55, 0.6, 0.6, size=22, color=ACCENT, align=PP_ALIGN.CENTER)

for i, (ic, title, desc) in enumerate(solutions):
    bx = box(s, 5.2, 1.4 + i*1.55, 4.0, 1.3, fill=RGBColor(0x05, 0x2A, 0x1A), radius=True)
    # 그린 좌측 바
    gb = s.shapes.add_shape(1, Inches(5.2), Inches(1.4 + i*1.55), Inches(0.06), Inches(1.3))
    gb.fill.solid(); gb.fill.fore_color.rgb = ACCENT; gb.line.fill.background()
    txt(s, ic,    5.4,  1.45 + i*1.55, 0.8, 0.8, size=28)
    txt(s, title, 6.1,  1.45 + i*1.55, 3.0, 0.45, size=16, bold=True, color=ACCENT2)
    txt(s, desc,  6.1,  1.85 + i*1.55, 2.9, 0.45, size=12, color=MUTED)

txt(s, "Before", 1.5, 1.05, 2, 0.35, size=13, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "After",  6.5, 1.05, 2, 0.35, size=13, color=ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# Slide 3 — 핵심 기능 3가지
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "핵심 기능", 0.4, 0.1, 6, 0.5, size=14, color=ACCENT, bold=True)
txt(s, "3가지 핵심 경쟁력", 0.4, 0.45, 10, 0.7, size=28, bold=True)

features = [
    ("🤖", "AI 효율 점수", ACCENT,
     ["GC%, Tm, 3'-end 안정성 등 8가지 특성 자동 분석",
      "점수 근거 투명 공개 — 성분별 기여도 표시",
      "XGBoost 학습 모델 준비 완료 (데이터 축적 시 즉시 전환)"]),
    ("🧪", "9단계 자동 파이프라인", RGBColor(0x60, 0xA5, 0xFA),
     ["MAFFT 정렬 → 보존 영역 탐색 → Primer3 후보 생성",
      "특이성 필터 → 변이 커버리지 → 열역학 검증",
      "AI 점수 → 순위화 → HTML/PDF 보고서 자동 생성"]),
    ("📊", "즉시 활용 보고서", RGBColor(0xF4, 0x72, 0x7E),
     ["상위 후보군 비교표, 서열 정보 포함",
      "AI 점수 근거 세부 내역 (클릭 확장)",
      "PDF 다운로드 & 공유 가능"]),
]

for i, (icon, title, color, bullets) in enumerate(features):
    bx = box(s, 0.3 + i*4.35, 1.35, 4.1, 5.5, fill=LIGHT_BG, radius=True)
    # 상단 색상 바
    cb = s.shapes.add_shape(1, Inches(0.3 + i*4.35), Inches(1.35), Inches(4.1), Inches(0.08))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()

    txt(s, icon,  0.5 + i*4.35, 1.5,  0.8, 0.9, size=34)
    txt(s, title, 1.3 + i*4.35, 1.55, 2.9, 0.8, size=17, bold=True, color=color)

    for j, bullet in enumerate(bullets):
        txt(s, f"• {bullet}", 0.55 + i*4.35, 2.55 + j*0.95, 3.8, 0.85,
            size=12.5, color=MUTED if j > 0 else WHITE, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 4 — 성과 지표 (숫자로 말한다)
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "상반기 성과", 0.4, 0.1, 6, 0.5, size=14, color=ACCENT, bold=True)
txt(s, "숫자로 말하는 성과", 0.4, 0.45, 10, 0.7, size=28, bold=True)

metrics = [
    ("9단계", "완전 자동화\n파이프라인", ACCENT),
    ("< 30초", "전체 파이프라인\n처리 시간", YELLOW),
    ("AI 점수\n8개 특성", "투명한 근거\n제공", RGBColor(0x60, 0xA5, 0xFA)),
    ("100%\n클라우드", "서버 없이\n운영 가능", RGBColor(0xF4, 0x72, 0x7E)),
]

for i, (num, label, color) in enumerate(metrics):
    bx = box(s, 0.3 + i*3.2, 1.5, 3.0, 2.5, fill=LIGHT_BG, radius=True)
    txt(s, num,   0.4 + i*3.2, 1.65, 2.8, 1.1, size=26, bold=True,
        color=color, align=PP_ALIGN.CENTER)
    txt(s, label, 0.4 + i*3.2, 2.75, 2.8, 0.9, size=13, color=MUTED,
        align=PP_ALIGN.CENTER)

# 타임라인 성과
txt(s, "주요 개발 마일스톤", 0.4, 4.25, 10, 0.4, size=16, bold=True)

milestones = [
    ("Phase 1", "핵심 파이프라인\n9단계 구현 완료", ACCENT),
    ("Phase 2", "Vercel 서버리스\n배포 성공", RGBColor(0x60, 0xA5, 0xFA)),
    ("Phase 3", "AI 점수 투명성\n보고서 고도화", YELLOW),
    ("진행 중", "XGBoost 모델\n데이터 수집", MUTED),
]

for i, (phase, desc, color) in enumerate(milestones):
    bx = box(s, 0.3 + i*3.2, 4.75, 3.0, 2.1, fill=LIGHT_BG, radius=True)
    # 상단 점
    dot = s.shapes.add_shape(9, Inches(1.7 + i*3.2), Inches(4.65), Inches(0.2), Inches(0.2))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    txt(s, phase, 0.5 + i*3.2, 4.85, 2.8, 0.45, size=13, bold=True,
        color=color, align=PP_ALIGN.CENTER)
    txt(s, desc,  0.4 + i*3.2, 5.3,  2.9, 0.9,  size=12, color=WHITE,
        align=PP_ALIGN.CENTER)

# 가로 연결선
line = s.shapes.add_shape(1, Inches(0.3), Inches(4.72), Inches(12.7), Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb = MUTED; line.line.fill.background()


# ════════════════════════════════════════════════════════════════
# Slide 5 — 기술 아키텍처 (간략)
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "기술 인프라", 0.4, 0.1, 6, 0.5, size=14, color=ACCENT, bold=True)
txt(s, "운영 비용 최소화 — 클라우드 네이티브 아키텍처", 0.4, 0.45, 12, 0.7, size=28, bold=True)

# 3-layer 다이어그램
layers = [
    ("사용자", ["웹 브라우저", "FASTA 업로드", "보고서 다운로드"],
     RGBColor(0x6E, 0xE7, 0xB7), 0.3),
    ("플랫폼 (Vercel)", ["FastAPI 백엔드", "9단계 파이프라인", "PDF/HTML 생성"],
     ACCENT, 4.5),
    ("데이터 (Supabase)", ["작업 상태 관리", "결과 영구 저장", "실시간 진행 표시"],
     RGBColor(0x60, 0xA5, 0xFA), 8.7),
]

for i, (title, items, color, x) in enumerate(layers):
    bx = box(s, x, 1.4, 3.8, 4.8, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(x), Inches(1.4), Inches(3.8), Inches(0.1))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()
    txt(s, title, x + 0.15, 1.55, 3.5, 0.55, size=16, bold=True, color=color)
    for j, item in enumerate(items):
        txt(s, f"  ✦  {item}", x + 0.2, 2.25 + j*0.85, 3.4, 0.75,
            size=13, color=WHITE if j == 0 else MUTED)

    if i < 2:
        txt(s, "⟷", x + 3.9, 3.5, 0.7, 0.6, size=24, color=ACCENT, align=PP_ALIGN.CENTER)

txt(s, "• 별도 서버 운용 불필요  • 자동 스케일링  • 99.9% 가용성  • 월정액 예측 비용",
    0.3, 6.5, 12.7, 0.5, size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# Slide 6 — 차별화 & 경쟁 우위
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "경쟁 우위", 0.4, 0.1, 6, 0.5, size=14, color=ACCENT, bold=True)
txt(s, "타 솔루션과 무엇이 다른가", 0.4, 0.45, 10, 0.7, size=28, bold=True)

headers = ["", "수작업\n전문가", "범용\n바이오인포", "OpenBioShield"]
col_colors = [DARK_BG, DARK_BG, DARK_BG, RGBColor(0x05, 0x2A, 0x1A)]
rows = [
    ("설계 소요 시간",     "수일~수주",  "수시간",     "< 30초"),
    ("AI 효율 예측",       "없음",       "제한적",     "✅ 8특성 분석"),
    ("결과 근거 투명성",   "없음",       "없음",       "✅ 성분별 분해"),
    ("변이 커버리지",      "수동 확인",  "별도 도구",  "✅ 자동 계산"),
    ("보고서 자동 생성",   "없음",       "없음",       "✅ HTML+PDF"),
    ("클라우드 즉시 접근", "없음",       "설치 필요",  "✅ 브라우저만"),
]

col_w = [3.0, 2.6, 2.6, 2.9]
col_x = [0.3, 3.4, 6.1, 9.0]

for ci, (header, cx) in enumerate(zip(headers, col_x)):
    fill = RGBColor(0x05, 0x2A, 0x1A) if ci == 3 else LIGHT_BG
    bx = box(s, cx, 1.35, col_w[ci], 0.7, fill=fill, radius=False)
    color = ACCENT if ci == 3 else (MUTED if ci == 0 else WHITE)
    txt(s, header, cx + 0.1, 1.4, col_w[ci] - 0.1, 0.6,
        size=13, bold=(ci == 3), color=color, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    for ci, (cell, cx) in enumerate(zip(row, col_x)):
        fill = RGBColor(0x05, 0x3A, 0x20) if ci == 3 else (LIGHT_BG if ri % 2 == 0 else RGBColor(0x16, 0x22, 0x34))
        bx = box(s, cx, 2.15 + ri * 0.75, col_w[ci], 0.73, fill=fill)
        color = ACCENT2 if (ci == 3 and "✅" in cell) else (WHITE if ci == 0 else MUTED)
        align = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        txt(s, cell, cx + 0.1, 2.2 + ri * 0.75, col_w[ci] - 0.1, 0.65,
            size=12.5, color=color, align=align, bold=(ci == 3))


# ════════════════════════════════════════════════════════════════
# Slide 7 — 로드맵 & 다음 단계
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "향후 계획", 0.4, 0.1, 6, 0.5, size=14, color=ACCENT, bold=True)
txt(s, "2026년 하반기 로드맵", 0.4, 0.45, 10, 0.7, size=28, bold=True)

roadmap = [
    ("Q3 2026", [
        "XGBoost 모델 학습 데이터 수집",
        "Multiplex qPCR 프로브 설계 고도화",
        "사용자 피드백 루프 구축",
    ], ACCENT, "완료 목표"),
    ("Q4 2026", [
        "외부 검증 데이터셋 연동",
        "LIMS 시스템 API 연동",
        "다국어 보고서 지원",
    ], YELLOW, "개발 예정"),
    ("2027+", [
        "규제 제출용 보고서 자동화 (IVD)",
        "멀티 타겟 동시 설계",
        "SaaS 외부 서비스화",
    ], RGBColor(0x60, 0xA5, 0xFA), "중장기 목표"),
]

for i, (quarter, items, color, badge) in enumerate(roadmap):
    bx = box(s, 0.3 + i*4.35, 1.4, 4.1, 5.5, fill=LIGHT_BG, radius=True)
    cb = s.shapes.add_shape(1, Inches(0.3 + i*4.35), Inches(1.4), Inches(4.1), Inches(0.08))
    cb.fill.solid(); cb.fill.fore_color.rgb = color; cb.line.fill.background()

    txt(s, quarter, 0.5 + i*4.35, 1.5,  3.7, 0.55, size=20, bold=True, color=color)
    txt(s, badge,   0.5 + i*4.35, 2.1,  3.7, 0.4,  size=11, color=MUTED)

    for j, item in enumerate(items):
        txt(s, f"→  {item}", 0.5 + i*4.35, 2.65 + j*1.0, 3.7, 0.85,
            size=13, color=WHITE, wrap=True)


# ════════════════════════════════════════════════════════════════
# Slide 8 — 마무리 (Call to Action)
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

bar2 = s.shapes.add_shape(1, 0, 0, Inches(0.12), prs.slide_height)
bar2.fill.solid(); bar2.fill.fore_color.rgb = ACCENT; bar2.line.fill.background()

txt(s, "OpenBioShield", 0.35, 1.8, 12, 1.0, size=46, bold=True, align=PP_ALIGN.CENTER)
txt(s, "분자진단 어세이 설계의 미래", 0.35, 2.95, 12.6, 0.7,
    size=26, color=ACCENT2, align=PP_ALIGN.CENTER)

# 핵심 메시지 3개
msgs = ["수일 → 30초", "AI 근거 투명 공개", "클라우드 즉시 배포"]
for i, msg in enumerate(msgs):
    bx = box(s, 1.3 + i*3.7, 3.9, 3.3, 0.75, fill=LIGHT_BG, radius=True)
    txt(s, msg, 1.4 + i*3.7, 3.95, 3.1, 0.65, size=16, bold=True,
        color=ACCENT, align=PP_ALIGN.CENTER)

txt(s, "문의 · 데모 요청 · 협업 제안", 0.35, 5.3, 12.6, 0.5,
    size=15, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "kyle.yoon.jh@gmail.com", 0.35, 5.75, 12.6, 0.5,
    size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── 저장 ──────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "OpenBioShield_성과보고_2026상반기.pptx")
prs.save(out)
print(f"저장 완료: {out}")
