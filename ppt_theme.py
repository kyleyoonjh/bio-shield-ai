"""
OpenBioShield PPT 공통 테마 — 씨젠 e-Sign 디자인 시스템 기반
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 ────────────────────────────────────────────────────
BG_LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)   # 콘텐츠 슬라이드 배경
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # 섹션 구분 배경
HDR_DARK  = RGBColor(0x0F, 0x17, 0x2A)   # 헤더 바
BLUE      = RGBColor(0x25, 0x63, 0xEB)   # 주 강조색
CYAN      = RGBColor(0x06, 0xB6, 0xD4)   # 보조 강조색
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x0F, 0x17, 0x2A)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
CARD_DARK = RGBColor(0x16, 0x23, 0x3C)
DIVIDER   = RGBColor(0x2D, 0x3F, 0x5E)
INFO_BG   = RGBColor(0xEF, 0xF6, 0xFF)
SUCCESS   = RGBColor(0x16, 0xA3, 0x4A)
SUCCESS_BG= RGBColor(0xF0, 0xFD, 0xF4)
WARN      = RGBColor(0xD9, 0x77, 0x06)
WARN_BG   = RGBColor(0xFF, 0xFB, 0xEB)
SLATE     = RGBColor(0x47, 0x55, 0x69)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_BG   = RGBColor(0xF1, 0xF5, 0xF9)
RED_DOT   = RGBColor(0xFC, 0xA5, 0xA5)
YLW_DOT   = RGBColor(0xFD, 0xE6, 0x8A)
GRN_DOT   = RGBColor(0x86, 0xEF, 0xAC)


def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill, line_fill=None):
    sh = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_fill:
        sh.line.color.rgb = line_fill
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh


def textbox(slide, text, x, y, w, h,
            size=13, bold=False, color=TEXT_DARK,
            align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


# ── 슬라이드 타입별 빌더 ────────────────────────────────────

def make_title_slide(prs, phase_num, phase_label, title, subtitle, accent=BLUE):
    """
    섹션 구분형 타이틀 슬라이드 (다크 배경)
    """
    s = add_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, BG_DARK)
    # 좌측 파란 바
    rect(s, 0, 0, 0.1, 7.5, accent)
    # 배경 대형 Phase 번호 (워터마크)
    textbox(s, phase_num, 7.5, 0.2, 6.0, 5.5,
            size=180, bold=True, color=RGBColor(0x1E, 0x29, 0x3B), align=PP_ALIGN.LEFT)
    # Phase 레이블
    textbox(s, phase_label, 0.4, 1.8, 7.0, 0.6,
            size=16, bold=True, color=accent)
    # 메인 타이틀
    textbox(s, title, 0.4, 2.45, 8.0, 1.4,
            size=42, bold=True, color=WHITE)
    # 서브타이틀
    textbox(s, subtitle, 0.4, 4.0, 9.0, 0.6,
            size=18, color=MUTED)
    # 하단 바
    rect(s, 0.4, 5.0, 3.5, 0.04, accent)
    return s


def make_section_slide(prs, section_num, section_title, section_desc, accent=BLUE):
    """
    섹션 인트로 슬라이드 (다크 배경, 번호 워터마크)
    """
    s = add_slide(prs)
    rect(s, 0, 0, 13.33, 7.5, BG_DARK)
    rect(s, 0, 0, 0.08, 7.5, accent)
    # 배경 워터마크 숫자
    textbox(s, section_num, 8.5, 0.3, 5.0, 5.0,
            size=200, bold=True, color=RGBColor(0x1E, 0x29, 0x3B))
    # 섹션 번호 크게
    textbox(s, section_num, 0.4, 2.2, 3.0, 1.3,
            size=72, bold=True, color=accent)
    # 섹션 제목
    textbox(s, section_title, 0.4, 3.7, 8.5, 1.2,
            size=36, bold=True, color=WHITE)
    # 설명
    textbox(s, section_desc, 0.4, 5.05, 9.0, 0.7,
            size=16, color=MUTED)
    return s


def make_content_header(slide, step_num, title, subtitle, accent=BLUE):
    """
    콘텐츠 슬라이드 공통 헤더 (다크 1.05" 헤더 바)
    """
    # 헤더 배경
    rect(slide, 0, 0, 13.33, 1.05, HDR_DARK)
    # 좌측 파란 바
    rect(slide, 0, 0, 0.07, 1.05, accent)
    # 스텝 번호 뱃지
    rect(slide, 0.28, 0.2, 0.62, 0.62, accent)
    textbox(slide, step_num, 0.28, 0.2, 0.62, 0.62,
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 타이틀
    textbox(slide, title, 1.05, 0.18, 11.5, 0.52,
            size=22, bold=True, color=WHITE)
    # 서브타이틀
    textbox(slide, subtitle, 1.05, 0.72, 11.5, 0.3,
            size=11, color=MUTED)


def step_card(slide, x, y, w, h, num, num_color, title, desc, title_color=TEXT_DARK):
    """
    번호 뱃지 + 설명 텍스트 카드 (흰 배경)
    """
    rect(slide, x, y, w, h, WHITE, BORDER)
    # 번호 뱃지
    rect(slide, x+0.12, y+(h-0.71)/2, 0.71, 0.71, num_color)
    textbox(slide, num, x+0.12, y+(h-0.71)/2, 0.71, 0.71,
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 제목 + 설명
    textbox(slide, title, x+1.0, y+0.12, w-1.15, 0.38,
            size=13, bold=True, color=title_color)
    if desc:
        textbox(slide, desc, x+1.0, y+0.52, w-1.15, h-0.65,
                size=11, color=SLATE, wrap=True)


def info_box(slide, x, y, w, h, label, content, accent=BLUE):
    """
    파란 왼쪽 선 인포 박스 (#EFF6FF 배경)
    """
    rect(slide, x, y, w, h, INFO_BG)
    rect(slide, x, y, 0.06, h, accent)
    if label:
        textbox(slide, label, x+0.2, y+0.12, w-0.3, 0.32,
                size=11, bold=True, color=accent)
    textbox(slide, content, x+0.2, y+(0.38 if label else 0.1),
            w-0.3, h-(0.38 if label else 0.1)-0.08,
            size=12, color=TEXT_DARK, wrap=True)


def success_box(slide, x, y, w, h, content):
    rect(slide, x, y, w, h, SUCCESS_BG)
    rect(slide, x, y, 0.06, h, SUCCESS)
    textbox(slide, content, x+0.2, y+0.1, w-0.3, h-0.15,
            size=12, color=SUCCESS, wrap=True)


def warn_box(slide, x, y, w, h, content):
    rect(slide, x, y, w, h, WARN_BG)
    rect(slide, x, y, 0.06, h, WARN)
    textbox(slide, content, x+0.2, y+0.1, w-0.3, h-0.15,
            size=12, color=WARN, wrap=True)


def dark_card_3col(slide, x, y, w, h, num, title, bullets, accent=CYAN):
    """
    다크 카드 (3열 레이아웃용) — 씨젠 레퍼런스 슬라이드2 스타일
    """
    rect(slide, x, y, w, h, CARD_DARK)
    rect(slide, x, y, w, 0.02, accent)          # 상단 선
    textbox(slide, num, x+0.07, y+0.08, 0.5, 0.42,
            size=13, bold=True, color=BLUE)
    textbox(slide, title, x+0.55, y+0.08, w-0.62, 0.42,
            size=12, bold=True, color=WHITE)
    rect(slide, x+0.07, y+0.56, w-0.14, 0.02, DIVIDER)
    for i, b in enumerate(bullets):
        textbox(slide, b, x+0.07, y+0.65+i*0.55, w-0.14, 0.5,
                size=9, bold=True, color=accent, wrap=True)


def label_value_row(slide, x, y, w, label, value, label_color=BLUE):
    rect(slide, x, y, 1.55, 0.72, label_color)
    textbox(slide, label, x, y, 1.55, 0.72,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, x+1.55, y, w-1.55, 0.72, WHITE, BORDER)
    textbox(slide, value, x+1.65, y+0.12, w-1.75, 0.48,
            size=13, color=TEXT_DARK, wrap=True)


def browser_mockup(slide, x, y, w, h, url=""):
    rect(slide, x, y, w, h, WHITE, BORDER)
    rect(slide, x, y, w, 0.35, GRAY_BG)
    # traffic lights
    rect(slide, x+0.15, y+0.085, 0.18, 0.18, RED_DOT)
    rect(slide, x+0.42, y+0.085, 0.18, 0.18, YLW_DOT)
    rect(slide, x+0.69, y+0.085, 0.18, 0.18, GRN_DOT)
    if url:
        textbox(slide, url, x+1.0, y+0.08, w-1.2, 0.22,
                size=10, color=SLATE)


def footer_bar(slide, x, y, w, content, accent=BLUE):
    rect(slide, x, y, w, 0.38, INFO_BG)
    textbox(slide, content, x+0.25, y+0.04, w-0.3, 0.3,
            size=12, color=accent)


def risk_badge_inline(slide, x, y, label, color):
    rect(slide, x, y, 1.25, 0.4, color)
    textbox(slide, label, x, y, 1.25, 0.4,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
