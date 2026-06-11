"""
OpenBioShield Phase 2 — 연구원 가이드 PPT
통계 검증 & 위험도 분석
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

DARK_BG  = RGBColor(0x0F, 0x17, 0x2A)
ACCENT   = RGBColor(0x10, 0xB9, 0x81)
ACCENT2  = RGBColor(0x6E, 0xE7, 0xB7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0x1E, 0x2D, 0x40)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)
BLUE     = RGBColor(0x60, 0xA5, 0xFA)
PINK     = RGBColor(0xF4, 0x72, 0x7E)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)
PURPLE   = RGBColor(0xC0, 0x84, 0xFC)
CARD2    = RGBColor(0x16, 0x22, 0x34)
VIOLET   = RGBColor(0x81, 0x8C, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def slide():   return prs.slides.add_slide(blank)

def bg(s, c=DARK_BG):
    sh = s.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()

def box(s, x, y, w, h, fill=LIGHT_BG, radius=False):
    sh = s.shapes.add_shape(17 if radius else 1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background()
    if radius: sh.adjustments[0]=0.05
    return sh

def top_bar(s, color=ACCENT):
    sh=s.shapes.add_shape(1,0,Inches(0.17),Inches(1.6),Inches(0.065))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background()

def side_bar(s, color=ACCENT):
    sh=s.shapes.add_shape(1,0,0,Inches(0.12),prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background()

def txt(s, text, x, y, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    txb.word_wrap=wrap
    tf=txb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    run=p.add_run(); run.text=text
    run.font.size=Pt(size); run.font.bold=bold
    run.font.color.rgb=color; run.font.italic=italic

def col_bar(s, x, y, h, color):
    sh=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(0.065),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background()

def tip_box(s, x, y, w, content, color=YELLOW):
    box(s, x, y, w, 0.72, fill=RGBColor(0x2A,0x20,0x05), radius=True)
    col_bar(s, x, y, 0.72, color)
    txt(s, f"💡  {content}", x+0.18, y+0.1, w-0.25, 0.52, size=12, color=YELLOW, wrap=True)

def risk_badge(s, x, y, level, color):
    b=s.shapes.add_shape(17,Inches(x),Inches(y),Inches(1.3),Inches(0.42))
    b.fill.solid(); b.fill.fore_color.rgb=color; b.line.fill.background()
    b.adjustments[0]=0.5
    txt(s, level, x+0.05, y+0.05, 1.2, 0.32, size=13, bold=True,
        align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# Slide 1 — 타이틀
# ════════════════════════════════════════════════════════════════
s=slide(); bg(s); side_bar(s, VIOLET)
txt(s, "Phase 2", 0.35, 1.1, 10, 0.55, size=16, color=VIOLET, bold=True)
txt(s, "통계 검증 &\n위험도 분석", 0.35, 1.6, 11, 1.8, size=52, bold=True)
txt(s, "CLSI EP05/EP09 · PFPS 위험도 엔진 · 데이터 수집", 0.35, 3.55, 12, 0.6, size=20, color=ACCENT2)
txt(s, "시약 연구원을 위한 기능 사용 가이드", 0.35, 4.25, 10, 0.5, size=15, color=MUTED)

features=[("📊","CLSI 통계 검증"),("⚠️","PFPS 위험도"),("🗄️","데이터 수집")]
for i,(ic,lb) in enumerate(features):
    box(s, 9.6, 1.5+i*1.55, 3.4, 1.3, fill=LIGHT_BG, radius=True)
    col_bar(s, 9.6, 1.5+i*1.55, 1.3, VIOLET)
    txt(s, ic, 9.8, 1.6+i*1.55, 0.8, 0.9, size=30)
    txt(s, lb, 10.6, 1.72+i*1.55, 2.2, 0.6, size=15, color=ACCENT2, bold=True)


# ════════════════════════════════════════════════════════════════
# Slide 2 — Phase 2 개요
# ════════════════════════════════════════════════════════════════
s=slide(); bg(s); top_bar(s, VIOLET)
txt(s, "Phase 2 개요", 0.4, 0.1, 6, 0.38, size=13, color=VIOLET, bold=True)
txt(s, "3개 화면으로 구성된 검증 & 위험도 워크플로우", 0.4, 0.45, 12, 0.55, size=24, bold=True)

flows=[
    ("1","📊 CLSI 통계 검증",       BLUE,   "정밀도·재현성 공식 검증",
     ["EP05: 반복성·재현성 분산분석","EP09: 방법비교 Deming 회귀","AI 보고서 자동 생성","결과 DB 자동 저장"]),
    ("2","⚠️ PFPS 위험도 엔진",     PINK,   "프라이머-변이 충돌 분석",
     ["변이 위치 거리 기반 위험 점수","3단계 위험등급 (HIGH/MEDIUM/LOW)","GPT-4o 서술형 설명","왓슨-크릭 결합 시뮬레이션"]),
    ("3","🗄️ 데이터 수집 현황",     VIOLET, "AI 모델 훈련 준비",
     ["수집된 검증 레코드 통계","질환·가이드라인별 분포","Fine-tuning 준비도 지표","연구원 피드백 집계"]),
]
for i,(num,title,color,sub,bullets) in enumerate(flows):
    bx=box(s,0.3+i*4.35,1.3,4.1,5.65,fill=LIGHT_BG,radius=True)
    cb=s.shapes.add_shape(1,Inches(0.3+i*4.35),Inches(1.3),Inches(4.1),Inches(0.08))
    cb.fill.solid(); cb.fill.fore_color.rgb=color; cb.line.fill.background()
    circ=s.shapes.add_shape(9,Inches(0.4+i*4.35),Inches(1.42),Inches(0.42),Inches(0.42))
    circ.fill.solid(); circ.fill.fore_color.rgb=color; circ.line.fill.background()
    txt(s,num,0.41+i*4.35,1.44,0.4,0.38,size=14,bold=True,align=PP_ALIGN.CENTER)
    txt(s,title,0.9+i*4.35,1.43,3.2,0.45,size=15,bold=True,color=color)
    txt(s,sub,  0.9+i*4.35,1.9, 3.2,0.38,size=12,color=MUTED)
    for j,b in enumerate(bullets):
        txt(s,f"  ›  {b}",0.5+i*4.35,2.45+j*0.72,3.7,0.65,size=13,
            color=WHITE if j==0 else MUTED,wrap=True)
    if i<2:
        txt(s,"➜",4.5+i*4.35,4.0,0.3,0.5,size=22,color=color,align=PP_ALIGN.CENTER)

tip_box(s,0.3,7.05,12.7,
        "Phase 2는 기존 프라이머·키트를 검증하는 단계입니다. CLSI 가이드라인 준수 여부를 확인하고 위험도를 문서화합니다.")


# ════════════════════════════════════════════════════════════════
# Slide 3 — CLSI EP05 정밀도 검증
# ════════════════════════════════════════════════════════════════
s=slide(); bg(s); top_bar(s, BLUE)
txt(s,"화면 1  |  CLSI EP05 — 정밀도(반복성·재현성) 검증",0.4,0.1,12,0.38,size=13,color=BLUE,bold=True)
txt(s,"데이터 업로드 한 번으로 ANOVA·CV 자동 계산",0.4,0.45,12,0.55,size=24,bold=True)

# 왼쪽: 단계별 사용법
box(s,0.3,1.3,5.9,5.55,fill=LIGHT_BG,radius=True)
txt(s,"⚙️  EP05 분석 단계별 가이드",0.5,1.42,5.5,0.42,size=15,bold=True,color=ACCENT2)
steps=[
    ("1","Excel 파일 업로드",BLUE,
     "CLSI EP05 형식 데이터 파일 드래그&드롭\n(샘플×반복×재현 측정값 포함)"),
    ("2","컬럼 매핑 확인",BLUE,
     "AI가 자동 감지한 컬럼 할당 검토\n잘못된 경우 드롭다운으로 직접 수정"),
    ("3","분석 실행",BLUE,
     "EP05 버튼 클릭 → ANOVA, CV%, 그랜드 평균 자동 계산\n완료까지 약 5–10초"),
    ("4","AI 보고서 확인",BLUE,
     "한국어·영어 토글 가능한 AI 생성 마크다운 보고서\n판정 근거·개선 권장사항 포함"),
    ("5","DB 자동 저장",BLUE,
     "분석 결과가 Supabase에 자동 저장\nPhase 2 데이터 수집 화면에서 누적 확인"),
]
for i,(num,title,color,desc) in enumerate(steps):
    box(s,0.45,2.02+i*0.88,5.6,0.8,fill=CARD2,radius=True)
    circ=s.shapes.add_shape(9,Inches(0.55),Inches(2.1+i*0.88),Inches(0.32),Inches(0.32))
    circ.fill.solid(); circ.fill.fore_color.rgb=color; circ.line.fill.background()
    txt(s,num,0.56,2.12+i*0.88,0.3,0.28,size=11,bold=True,align=PP_ALIGN.CENTER)
    txt(s,title,0.98,2.1+i*0.88,2.5,0.38,size=13,bold=True,color=color)
    txt(s,desc, 0.98,2.5+i*0.88, 4.9,0.35,size=11.5,color=MUTED,wrap=True)

# 오른쪽: 결과 해석
box(s,6.45,1.3,6.55,2.8,fill=LIGHT_BG,radius=True)
txt(s,"📊  EP05 결과 수치 해석",6.65,1.42,6.1,0.42,size=15,bold=True,color=ACCENT2)
ep05_metrics=[
    ("반복성 CV%",    "< 5%",   "동일 기기 반복 측정 정밀도",    BLUE),
    ("재현성 CV%",    "< 10%",  "다일·다기기 간 재현성",          ACCENT),
    ("ANOVA F-값",    "p < 0.05","측정 그룹 간 유의한 차이 여부", YELLOW),
    ("그랜드 평균",   "기준값 ±5%","전체 측정값 중심 경향",        MUTED),
]
for i,(metric,ideal,desc,color) in enumerate(ep05_metrics):
    box(s,6.6,1.98+i*0.62,6.25,0.54,fill=CARD2,radius=True)
    col_bar(s,6.6,1.98+i*0.62,0.54,color)
    txt(s,metric,6.8,2.03+i*0.62,2.0,0.44,size=12,bold=True,color=color)
    txt(s,ideal, 8.85,2.03+i*0.62,1.5,0.44,size=12,color=WHITE,align=PP_ALIGN.CENTER)
    txt(s,desc,  10.4,2.03+i*0.62,2.5,0.44,size=11,color=MUTED,align=PP_ALIGN.CENTER)

# 오른쪽 하단: 판정 기준
box(s,6.45,4.3,6.55,2.55,fill=LIGHT_BG,radius=True)
txt(s,"✅  합격/불합격 판정 기준",6.65,4.42,6.1,0.42,size=14,bold=True,color=ACCENT2)
judgments=[
    ("반복성 CV < 5%  AND  재현성 CV < 10%","합격 — 다음 단계 진행",ACCENT),
    ("반복성 CV 5–10% OR  재현성 CV 10–15%","조건부 합격 — 추가 검증",YELLOW),
    ("CV > 15% 또는 ANOVA p > 0.05","불합격 — 프로토콜 재검토",PINK),
]
for i,(cond,verdict,color) in enumerate(judgments):
    box(s,6.6,4.98+i*0.62,6.25,0.55,fill=CARD2,radius=True)
    col_bar(s,6.6,4.98+i*0.62,0.55,color)
    txt(s,cond,   6.8,5.02+i*0.62,3.8,0.45,size=11,color=MUTED,wrap=True)
    txt(s,verdict,10.65,5.02+i*0.62,2.1,0.45,size=12,bold=True,color=color)

tip_box(s,0.3,7.05,12.7,
        "재현성 CV가 10%를 초과하면 PFPS 위험도 엔진에서 자동으로 HIGH 위험으로 분류됩니다. CV 결과를 반드시 위험도 분석과 연계하세요.")


# ════════════════════════════════════════════════════════════════
# Slide 4 — CLSI EP09 방법비교 검증
# ════════════════════════════════════════════════════════════════
s=slide(); bg(s); top_bar(s, TEAL:=RGBColor(0x2D,0xD4,0xBF))
txt(s,"화면 1  |  CLSI EP09 — 방법비교(Method Comparison) 검증",0.4,0.1,12,0.38,size=13,color=TEAL,bold=True)
txt(s,"신규 진단법 vs 기준법 — 동등성 통계 검증",0.4,0.45,12,0.55,size=24,bold=True)

# 왼쪽: 분석 항목
box(s,0.3,1.3,5.9,5.55,fill=LIGHT_BG,radius=True)
txt(s,"📐  EP09 분석 항목 및 해석",0.5,1.42,5.5,0.42,size=15,bold=True,color=ACCENT2)
ep09_items=[
    ("Deming 회귀",         TEAL,
     "기울기 1±0.1, 절편 0±10%\n→ 두 방법이 측정값 스케일이 동일한지 확인\n기울기가 1에서 멀수록 비례 오차 존재"),
    ("Bland-Altman 분석",   BLUE,
     "평균 차이(Bias) ±5% 이내\n→ 두 방법 간 체계적 차이 시각화\n95% 일치구간 확인"),
    ("Pearson 상관계수(r)", YELLOW,
     "r ≥ 0.95 합격 기준\n→ 두 방법 간 선형 관련성 강도\n0.90 미만이면 방법 재검토"),
]
for i,(title,color,desc) in enumerate(ep09_items):
    box(s,0.45,2.05+i*1.52,5.6,1.38,fill=CARD2,radius=True)
    col_bar(s,0.45,2.05+i*1.52,1.38,color)
    txt(s,title,0.65,2.12+i*1.52,5.1,0.45,size=14,bold=True,color=color)
    txt(s,desc, 0.65,2.6 +i*1.52,5.1,0.75,size=12,color=MUTED,wrap=True)

# 오른쪽: 활용 시나리오
box(s,6.45,1.3,6.55,5.55,fill=LIGHT_BG,radius=True)
txt(s,"🔬  실무 활용 시나리오",6.65,1.42,6.1,0.42,size=15,bold=True,color=ACCENT2)
ep09_cases=[
    ("신규 PCR 키트 허가 신청", TEAL,
     "기준 방법(PCR Gold Standard)과 신규 키트 동시 측정\n→ EP09 업로드 → r≥0.95, Bias±5% 확인\n→ 허가 서류용 AI 보고서 다운로드"),
    ("장비 교체 동등성 확인", BLUE,
     "구형 장비 vs 신형 장비 동일 샘플 측정\n→ Deming 기울기 확인\n→ 기울기 1.0±0.1이면 교체 승인"),
    ("플랫폼 간 비교 (POC vs Lab)", YELLOW,
     "현장 검사 장비(POC)와 검사실 표준 장비 비교\n→ Bland-Altman Bias 확인\n→ 허용 오차 범위 내이면 POC 도입 승인"),
    ("시약 로트 변경 검증", ORANGE,
     "기존 로트 vs 신규 로트 비교\n→ 기울기·절편 범위 내 확인\n→ 품질관리 문서 자동 생성"),
]
for i,(title,color,desc) in enumerate(ep09_cases):
    box(s,6.6,2.0+i*1.15,6.25,1.02,fill=CARD2,radius=True)
    col_bar(s,6.6,2.0+i*1.15,1.02,color)
    txt(s,title,6.8,2.07+i*1.15,5.8,0.42,size=13,bold=True,color=color)
    txt(s,desc, 6.8,2.52+i*1.15,5.8,0.44,size=11.5,color=MUTED,wrap=True)

tip_box(s,0.3,7.05,12.7,
        "AI 보고서는 한국어·영어 전환이 가능합니다. 허가 신청 시 영어 버전을 출력하면 해외 규제 제출 자료로 바로 활용할 수 있습니다.")


# ════════════════════════════════════════════════════════════════
# Slide 5 — PFPS 위험도 엔진
# ════════════════════════════════════════════════════════════════
TEAL=RGBColor(0x2D,0xD4,0xBF)
s=slide(); bg(s); top_bar(s, PINK)
txt(s,"화면 2  |  PFPS 위험도 엔진",0.4,0.1,12,0.38,size=13,color=PINK,bold=True)
txt(s,"프라이머-변이 충돌 위험을 점수로 수치화",0.4,0.45,12,0.55,size=24,bold=True)

# 왼쪽: PFPS 개념 설명
box(s,0.3,1.3,5.9,5.55,fill=LIGHT_BG,radius=True)
txt(s,"⚠️  PFPS란?",0.5,1.42,5.5,0.42,size=15,bold=True,color=ACCENT2)
txt(s,
    "Primer Fitness Position Score\n\n"
    "변이 위치가 프라이머 결합 구간과 얼마나 가까운지를\n"
    "거리 가중치로 계산한 위험 점수입니다.\n\n"
    "3'-end(3 프라임 말단)에 가까울수록 PCR 실패 위험이\n"
    "기하급수적으로 높아집니다.",
    0.5,1.95,5.6,1.7,size=13,color=WHITE)

# 위치별 위험도
risk_zones=[
    ("0–3 bp", "3'-end 내부",    "HIGH",   PINK,   "증폭 실패 확률 매우 높음\n즉각적 재설계 필요"),
    ("4–15 bp","결합 중간 부위", "MEDIUM", ORANGE, "증폭 효율 저하 가능\nTm 저하, Gradient PCR 고려"),
    ("16+ bp", "결합 외곽",      "LOW",    ACCENT, "영향 미미\n모니터링 수준 유지"),
]
for i,(dist,zone,level,color,desc) in enumerate(risk_zones):
    box(s,0.45,3.85+i*1.0,5.6,0.9,fill=CARD2,radius=True)
    col_bar(s,0.45,3.85+i*1.0,0.9,color)
    txt(s,dist, 0.65,3.9 +i*1.0,1.0,0.42,size=13,bold=True,color=color)
    txt(s,zone, 1.7, 3.9 +i*1.0,1.5,0.42,size=12,color=MUTED)
    risk_badge(s,3.3,3.9+i*1.0,level,color)
    txt(s,desc, 0.65,4.35+i*1.0,5.1,0.38,size=11,color=MUTED,wrap=True)

# 오른쪽: 결과 화면 구성 & 활용
box(s,6.45,1.3,6.55,2.65,fill=LIGHT_BG,radius=True)
txt(s,"📋  결과 화면 구성 요소",6.65,1.42,6.1,0.42,size=15,bold=True,color=ACCENT2)
pfps_ui=[
    ("위험도 등급 카드",      PINK,   "HIGH/MEDIUM/LOW 색상 배지 + 수치"),
    ("결합 시뮬레이션 그림",  ORANGE, "프라이머↔타겟 서열 정렬 표시"),
    ("GPT-4o 설명",          VIOLET, "위험 원인·개선 방향 서술형 해설"),
    ("변이 위치 목록",        BLUE,   "각 미스매치의 bp 거리·영향도"),
]
for i,(name,color,desc) in enumerate(pfps_ui):
    box(s,6.6,1.98+i*0.52,6.25,0.46,fill=CARD2,radius=True)
    col_bar(s,6.6,1.98+i*0.52,0.46,color)
    txt(s,name,6.8, 2.02+i*0.52,2.4,0.38,size=12,bold=True,color=color)
    txt(s,desc,9.25,2.02+i*0.52,3.6,0.38,size=11.5,color=MUTED)

# 오른쪽 하단: 판단 기준
box(s,6.45,4.15,6.55,2.7,fill=LIGHT_BG,radius=True)
txt(s,"🔬  PFPS 결과 활용 판단 기준",6.65,4.27,6.1,0.42,size=14,bold=True,color=ACCENT2)
pfps_actions=[
    ("HIGH",  PINK,   "즉각 재설계 → Phase 3으로 신규 프라이머 설계"),
    ("MEDIUM",ORANGE, "실험 진행 + Gradient PCR 병행 → 최적화 검증"),
    ("LOW",   ACCENT, "사용 적합 — 분기별 변이 업데이트 확인 권장"),
]
for i,(level,color,action) in enumerate(pfps_actions):
    box(s,6.6,4.88+i*0.62,6.25,0.55,fill=CARD2,radius=True)
    col_bar(s,6.6,4.88+i*0.62,0.55,color)
    risk_badge(s,6.7,4.93+i*0.62,level,color)
    txt(s,action,8.1,4.93+i*0.62,4.7,0.45,size=12,color=MUTED,wrap=True)

tip_box(s,0.3,7.05,12.7,
        "재현성 CV > 10% 이면 PFPS 엔진이 HIGH로 자동 에스컬레이션합니다. CLSI EP05 결과와 PFPS를 반드시 함께 확인하세요.")


# ════════════════════════════════════════════════════════════════
# Slide 6 — 데이터 수집 & AI 모델 발전
# ════════════════════════════════════════════════════════════════
s=slide(); bg(s); top_bar(s, VIOLET)
txt(s,"화면 3  |  데이터 수집 현황 대시보드",0.4,0.1,12,0.38,size=13,color=VIOLET,bold=True)
txt(s,"검증 데이터 축적 → AI 모델 자동 개선",0.4,0.45,12,0.55,size=24,bold=True)

# 왼쪽: 화면 구성
box(s,0.3,1.3,5.9,5.55,fill=LIGHT_BG,radius=True)
txt(s,"📊  화면 구성 요소",0.5,1.42,5.5,0.42,size=15,bold=True,color=ACCENT2)
dc_items=[
    ("수집 레코드 수", VIOLET,
     "EP05·EP09 분석 완료 시 자동 저장된\n총 레코드 수 — 목표: 5,000건 이상"),
    ("질환별 분포",    BLUE,
     "SARS-CoV-2·HPV·STI 각 질환별\n데이터 비율 도넛 차트"),
    ("위험도 분포",    PINK,
     "HIGH/MEDIUM/LOW 비율\n치우침이 클 경우 데이터 균형화 필요"),
    ("피드백 정확도",  YELLOW,
     "연구원 피드백(accurate/partially/incorrect)\n누적 집계 — AI 보고서 품질 지표"),
    ("Fine-tuning 준비도", ACCENT,
     "5,000건 도달률 진행 바\n목표 달성 시 XGBoost 재학습 자동 알림"),
]
for i,(title,color,desc) in enumerate(dc_items):
    box(s,0.45,2.02+i*0.9,5.6,0.8,fill=CARD2,radius=True)
    col_bar(s,0.45,2.02+i*0.9,0.8,color)
    txt(s,title,0.65,2.07+i*0.9,2.2,0.42,size=13,bold=True,color=color)
    txt(s,desc, 2.95,2.07+i*0.9,3.1,0.65,size=11.5,color=MUTED,wrap=True)

# 오른쪽: AI 발전 사이클
box(s,6.45,1.3,6.55,5.55,fill=LIGHT_BG,radius=True)
txt(s,"🔄  데이터 → AI 개선 사이클",6.65,1.42,6.1,0.42,size=15,bold=True,color=ACCENT2)

cycle=[
    ("1","EP05·EP09 분석 실행",      VIOLET, "시약 연구원이 검증 데이터 업로드"),
    ("2","결과 자동 저장",            BLUE,   "분석 결과가 Supabase DB에 적재"),
    ("3","피드백 입력",               YELLOW, "AI 보고서 정확도 평가 (3단계)"),
    ("4","데이터 5,000건 도달",       PINK,   "Fine-tuning 준비 완료 알림"),
    ("5","XGBoost 모델 재학습",       ACCENT, "Phase 3 AI 점수 정확도 향상"),
    ("6","더 정확한 자동 설계",       VIOLET, "연구원 수작업 검토 시간 감소"),
]
for i,(num,title,color,desc) in enumerate(cycle):
    box(s,6.6,2.0+i*0.75,6.25,0.65,fill=CARD2,radius=True)
    circ=s.shapes.add_shape(9,Inches(6.7),Inches(2.07+i*0.75),Inches(0.32),Inches(0.32))
    circ.fill.solid(); circ.fill.fore_color.rgb=color; circ.line.fill.background()
    txt(s,num,6.71,2.09+i*0.75,0.3,0.28,size=11,bold=True,align=PP_ALIGN.CENTER)
    txt(s,title,7.12,2.08+i*0.75,2.5,0.38,size=13,bold=True,color=color)
    txt(s,desc, 9.65,2.08+i*0.75,3.1,0.38,size=11.5,color=MUTED,wrap=True)
    if i<5:
        txt(s,"↓",9.7,2.69+i*0.75,0.3,0.08,size=11,color=MUTED,align=PP_ALIGN.CENTER)

tip_box(s,0.3,7.05,12.7,
        "분석 후 AI 보고서 피드백 버튼(accurate·partially·incorrect)을 반드시 클릭해주세요. 이 피드백이 AI 모델 개선의 핵심 데이터입니다.")


# ════════════════════════════════════════════════════════════════
# Slide 7 — 규제 대응 & 문서화 활용
# ════════════════════════════════════════════════════════════════
s=slide(); bg(s); top_bar(s, YELLOW)
txt(s,"규제 대응  |  Phase 2 결과의 문서화 활용",0.4,0.1,12,0.38,size=13,color=YELLOW,bold=True)
txt(s,"식약처·CE·FDA 허가 대응 자료로 바로 활용",0.4,0.45,12,0.55,size=24,bold=True)

# 규제별 활용
regs=[
    ("🇰🇷  식약처\n(MFDS)",  BLUE,   "IVD 허가",
     ["EP05: 정밀도 성능 자료 (GD-01)", "EP09: 방법비교 자료",
      "AI 보고서 국문 출력 → 직접 첨부", "PFPS 위험도 근거 문서"]),
    ("🇪🇺  CE-IVD\n(IVDR)",  VIOLET, "유럽 인증",
     ["EP05 CV% 표 → Technical Documentation",
      "EP09 Deming 회귀 그래프 → 영문 출력",
      "Statistical Summary 자동 생성", "Risk Management 파일 첨부"]),
    ("🇺🇸  FDA 510(k)\n(EUA)", PINK,  "미국 허가",
     ["Analytical Performance 섹션",
      "Method Comparison Data",
      "AI 보고서 영문 출력",
      "Reproducibility Study 요약"]),
]
for i,(title,color,badge_txt,items) in enumerate(regs):
    bx=box(s,0.3+i*4.35,1.3,4.1,5.55,fill=LIGHT_BG,radius=True)
    cb=s.shapes.add_shape(1,Inches(0.3+i*4.35),Inches(1.3),Inches(4.1),Inches(0.08))
    cb.fill.solid(); cb.fill.fore_color.rgb=color; cb.line.fill.background()
    txt(s,title,0.5+i*4.35,1.42,2.5,0.75,size=16,bold=True,color=color)
    b2=s.shapes.add_shape(17,Inches(2.8+i*4.35),Inches(1.45),Inches(1.4),Inches(0.4))
    b2.fill.solid(); b2.fill.fore_color.rgb=CARD2; b2.line.fill.background()
    txt(s,badge_txt,2.85+i*4.35,1.5,1.35,0.35,size=11,color=color,align=PP_ALIGN.CENTER)
    for j,item in enumerate(items):
        txt(s,f"  ›  {item}",0.5+i*4.35,2.42+j*0.85,3.7,0.75,
            size=13,color=WHITE if j==0 else MUTED,wrap=True)

tip_box(s,0.3,7.05,12.7,
        "규제 제출용 문서는 반드시 영문 AI 보고서를 사용하세요. 분석 완료 후 언어 토글 버튼으로 즉시 전환 가능합니다.")


# ════════════════════════════════════════════════════════════════
# Slide 8 — Phase 2 체크리스트
# ════════════════════════════════════════════════════════════════
s=slide(); bg(s); top_bar(s, ORANGE)
txt(s,"Phase 2 완료 체크리스트",0.4,0.1,12,0.38,size=13,color=ORANGE,bold=True)
txt(s,"다음 단계(Phase 3) 진행 전 반드시 확인",0.4,0.45,12,0.55,size=24,bold=True)

checks=[
    (BLUE,   "EP05 완료",       "반복성 CV < 5%,  재현성 CV < 10% 달성",
     "미달 시: 측정 조건 재검토 (온도, 검체 균질도) 후 재분석"),
    (TEAL,   "EP09 완료",       "r ≥ 0.95, Deming 기울기 1.0±0.1, Bias < 5%",
     "미달 시: 측정 범위·기준법 선정 재검토"),
    (PINK,   "PFPS 확인",       "HIGH 위험 프라이머 없음",
     "HIGH 존재 시: Phase 3으로 해당 타겟 재설계 진행"),
    (YELLOW, "AI 보고서 검토",  "보고서 내용 확인 후 피드백(accurate/partly/incorrect) 입력",
     "피드백 없으면 데이터 수집 카운트에서 제외됨"),
    (VIOLET, "DB 저장 확인",    "데이터 수집 화면에서 오늘 분석 레코드 수 증가 확인",
     "저장 실패 시: 네트워크 오류 — 재분석 또는 개발팀 문의"),
    (ORANGE, "문서 저장",       "규제 제출용 PDF·AI 보고서 다운로드 및 ELN 첨부",
     "허가 신청 전 보고서 버전 관리 필수"),
    (ACCENT, "Phase 3 준비",    "재설계 필요 타겟 목록 정리, FASTA 파일 구성",
     "Phase 1 DNA 맵에서 보존 영역 재확인 후 서열 수집"),
    (MUTED,  "팀 공유",         "분석 결과 공유 — 이메일 또는 Slack에 PDF 첨부",
     "보고서 링크 대신 PDF 파일 첨부 권장 (클라우드 접근 불가 환경 대비)"),
]
for i,(color,title,check,action) in enumerate(checks):
    row=i%4; col=i//4
    bx=box(s,0.3+col*6.6,1.35+row*1.42,6.3,1.3,fill=LIGHT_BG,radius=True)
    cb=s.shapes.add_shape(1,Inches(0.3+col*6.6),Inches(1.35+row*1.42),Inches(0.065),Inches(1.3))
    cb.fill.solid(); cb.fill.fore_color.rgb=color; cb.line.fill.background()
    chk=s.shapes.add_shape(1,Inches(0.5+col*6.6),Inches(1.5+row*1.42),Inches(0.3),Inches(0.3))
    chk.fill.solid(); chk.fill.fore_color.rgb=CARD2; chk.line.color.rgb=color
    txt(s,title,0.95+col*6.6,1.47+row*1.42,3.0,0.4,size=14,bold=True,color=color)
    txt(s,check,0.5+col*6.6, 1.9 +row*1.42,5.7,0.38,size=12,color=WHITE)
    txt(s,f"→ {action}",0.5+col*6.6,2.3+row*1.42,5.7,0.28,size=10.5,color=MUTED,wrap=True)


out=os.path.join(os.path.dirname(__file__),"OpenBioShield_Phase2_연구원가이드.pptx")
prs.save(out)
print(f"저장 완료: {out}")
